#!/usr/bin/env python3
"""
build_axisless_pptx.py — Build an axisless PowerPoint from generated axisless figures.

Creates Cardiac_RODEO_Axisless.pptx with the same layout as the tracked version,
but using axisless images (no axis labels, titles, ticks, spines).

Positions extracted from user-edited PPTX. Row heights enforced for consistency:
all panels in the same row share the same height, center-aligned vertically.

Falls back to original images for any panels without axisless versions.

Usage:
    python build_axisless_pptx.py
"""

from pptx import Presentation
from pptx.util import Inches, Emu, Pt
from pptx.enum.text import PP_ALIGN
from pathlib import Path
from PIL import Image
import json

# ============================================================================
# PATHS
# ============================================================================

SCRIPT_DIR = Path(__file__).parent.resolve()
PROJECT_ROOT = SCRIPT_DIR.parent.parent

FIGURES_DIR = SCRIPT_DIR
OUTPUT_PPTX = SCRIPT_DIR / "Cardiac_RODEO_Axisless.pptx"
LAYOUT_JSON = SCRIPT_DIR / "slide_layout.json"

# Slide dimensions (EMU) — matches the tracked presentation
SLIDE_WIDTH = 6483350   # ~7.09"
SLIDE_HEIGHT = 7745413  # ~8.47"
EMU_PER_INCH = 914400

# Row grouping tolerance for height enforcement
ROW_TOLERANCE = 0.4  # inches


# ============================================================================
# IMAGE RESOLUTION
# ============================================================================

def _find_image(fig_id, letter, filename=None):
    """Find the axisless image for a panel, falling back to the original."""
    fig_dir = FIGURES_DIR / f'Fig_{fig_id}'
    axisless_dir = fig_dir / 'Axisless'

    candidates = []
    if filename:
        candidates.append(filename)
    candidates.append(f'Fig_{fig_id}{letter}.png')

    # Try axisless first
    for name in candidates:
        p = axisless_dir / name
        if p.exists():
            return p, True

    # Fall back to original
    for name in candidates:
        p = fig_dir / name
        if p.exists():
            return p, False

    return None, False


# ============================================================================
# ROW HEIGHT ENFORCEMENT
# ============================================================================

def _enforce_row_heights(images, exclude_names=None):
    """
    Group images into rows by y position, then enforce:
    - All images in a row get the same height (max in row)

    ONLY changes height — x, y, w are preserved exactly as the user set them.

    images: list of (letter, filename, x, y, w, h) tuples
    exclude_names: set of filenames to skip (e.g. colorbars)
    Returns: new list with adjusted h values only.
    """
    if exclude_names is None:
        exclude_names = set()

    # Separate excluded items
    included = [(i, img) for i, img in enumerate(images) if img[1] not in exclude_names]

    # Group into rows by y position
    rows = []
    used = set()
    sorted_items = sorted(included, key=lambda x: x[1][3])  # sort by y

    for idx, img in sorted_items:
        if idx in used:
            continue
        row = [(idx, img)]
        used.add(idx)
        ref_center = img[3] + img[5] / 2  # y + h/2

        for idx2, img2 in sorted_items:
            if idx2 in used:
                continue
            center2 = img2[3] + img2[5] / 2
            if abs(center2 - ref_center) < ROW_TOLERANCE:
                row.append((idx2, img2))
                used.add(idx2)

        rows.append(row)

    # Only snap heights — preserve x, y, w exactly
    result = list(images)
    for row in rows:
        if len(row) < 2:
            continue
        max_h = max(img[5] for _, img in row)
        for idx, img in row:
            letter, fname, x, y, w, h = img
            result[idx] = (letter, fname, x, y, w, round(max_h, 3))

    return result


# ============================================================================
# SLIDE HELPERS
# ============================================================================

def _add_slide(prs, title_text=""):
    """Add a blank slide with optional title."""
    layout = prs.slide_layouts[6]  # Blank
    slide = prs.slides.add_slide(layout)

    if title_text:
        txBox = slide.shapes.add_textbox(
            Inches(0.2), Inches(0.1), Inches(6.5), Inches(0.5))
        tf = txBox.text_frame
        tf.word_wrap = True
        p = tf.paragraphs[0]
        p.text = title_text
        p.font.size = Pt(14)
        p.font.bold = True
        p.alignment = PP_ALIGN.LEFT

    return slide


def _place_image(slide, img_path, x_in, y_in, w_in, h_in):
    """Place an image on a slide at the given position (inches)."""
    if img_path is None or not Path(img_path).exists():
        return
    slide.shapes.add_picture(
        str(img_path),
        Emu(int(x_in * EMU_PER_INCH)),
        Emu(int(y_in * EMU_PER_INCH)),
        Emu(int(w_in * EMU_PER_INCH)),
        Emu(int(h_in * EMU_PER_INCH)),
    )


def build_panel_slide(prs, fig_id, panels, layout, title):
    """Build a slide with panels at positions from slide_layout.json."""
    slide = _add_slide(prs, title)
    placed = 0
    fallback = 0

    for item in panels:
        letter = item[0] if isinstance(item, tuple) else item
        filename = item[1] if isinstance(item, tuple) and len(item) > 1 else None

        key = f'Fig_{fig_id}{letter}'
        pos = layout.get(key)
        if pos is None:
            continue

        img_path, is_axisless = _find_image(fig_id, letter, filename)
        if img_path:
            _place_image(slide, img_path, pos['x'], pos['y'], pos['w'], pos['h'])
            placed += 1
            if not is_axisless:
                fallback += 1

    status = f"  Slide {fig_id}: {placed} panels"
    if fallback:
        status += f" ({fallback} fallback to original)"
    print(status)
    return slide


# ============================================================================
# SLIDE 2 — Extracted positions from user-edited PPTX (MD5-verified mapping)
# ============================================================================

SLIDE2_IMAGES = [
    # (letter, filename, x, y, w, h)
    # Row 1: SNR analysis
    ('d', 'Fig_2d.png',                                  0.288, 3.700, 2.300, 1.300),
    # Row 2: Dose-response + heatmap
    ('g', 'Fig_2g_Epirubicin_O2.png',                    0.597, 5.482, 1.811, 1.300),
    ('h', 'Fig_2h_Epirubicin_TC50.png',                  2.682, 5.159, 2.000, 1.300),
    ('i', 'Fig_2i_Epirubicin_O2_Heatmap.png',            4.792, 5.486, 2.326, 1.292),
    # Row 3: Mexiletine panels
    ('j', 'Fig_2j_Mexiletine_Contractility.png',          0.370, 6.980, 1.809, 1.330),
    ('k', 'Fig_2k_Mexiletine_Waveforms.png',              2.525, 6.976, 2.284, 1.387),
    ('l', 'Fig_2l_Mexiletine_Contractility_Heatmap.png',  4.788, 7.000, 2.284, 1.290),
]


# ============================================================================
# SLIDE 3 — Extracted positions from user-edited PPTX (MD5-verified mapping)
# ============================================================================

SLIDE3_IMAGES = [
    # (letter, filename, x, y, w, h)
    # Row 1: Heatmaps + 3D surfaces + colorbar
    ('a', 'Fig_3a_Dactinomycin_O2_Heatmap.png',          0.288, 1.098, 0.984, 0.984),
    ('a', 'Fig_3a_Mexiletine_O2_Heatmap.png',             4.510, 1.047, 0.984, 0.984),
    ('b', 'Mexiletine_Eq7_biphasic_response.png',          5.495, 0.917, 1.181, 1.181),
    ('b', 'Dactinomycin_Eq3_gaussian_hill_hybrid.png',     1.251, 0.998, 1.181, 1.181),
    ('b', 'Nifedipine_Eq10_modified_hill_simple.png',      3.188, 0.949, 1.181, 1.181),
    ('b', 'Fig_3b_colorbar.png',                           6.802, 0.998, 0.130, 1.000),
    # Row 2: R² chart + scatter plot
    ('c', 'Fig_3c.png',                                    0.504, 2.342, 1.163, 1.161),
    ('d', 'Fig_3d.png',                                    3.088, 2.342, 3.055, 1.161),
    # Row 3: Well vs model overlays
    ('e', 'Fig_3e_O2.png',                                 0.408, 3.661, 2.000, 1.366),
    ('e', 'Fig_3e_Contractility.png',                      2.633, 3.716, 2.000, 1.366),
]


def build_slide_2(prs):
    """Build slide 2 with row-height-enforced layout."""
    images = _enforce_row_heights(SLIDE2_IMAGES)
    slide = _add_slide(prs, "Figure 2")
    placed = 0
    for letter, filename, x, y, w, h in images:
        img_path, _ = _find_image('2', letter, filename)
        if img_path:
            _place_image(slide, img_path, x, y, w, h)
            placed += 1
    print(f"  Slide 2: {placed} panels (row heights enforced)")
    return slide


def build_slide_3(prs):
    """Build slide 3 with row-height-enforced compound panels."""
    images = _enforce_row_heights(SLIDE3_IMAGES, exclude_names={'Fig_3b_colorbar.png'})
    slide = _add_slide(prs, "Figure 3")
    placed = 0
    for letter, filename, x, y, w, h in images:
        img_path, _ = _find_image('3', letter, filename)
        if img_path:
            _place_image(slide, img_path, x, y, w, h)
            placed += 1
    print(f"  Slide 3: {placed} panels (row heights enforced)")
    return slide


# ============================================================================
# 5x5 GRID SLIDES (4 & 5)
# ============================================================================

def build_5x5_slide(prs, fig_num, resp_type, title):
    """Build a 5x5 grid slide using axisless individual images."""
    slide = _add_slide(prs, title)

    GRID_LEFT = 416424
    GRID_TOP = 533120
    H_STEP = 981300
    V_STEP = 979431
    CELL_SIZE = int(1.22 * EMU_PER_INCH)
    RIGHT_MARGIN = 137160

    axisless_dir = FIGURES_DIR / f'Fig_{fig_num}' / f'{resp_type}_5x5_Individual_Axisless'
    orig_dir = FIGURES_DIR / f'Fig_{fig_num}' / f'{resp_type}_5x5_Individual'

    img_dir = axisless_dir if axisless_dir.exists() else orig_dir
    source = "axisless" if axisless_dir.exists() else "original"

    if not img_dir.exists():
        print(f"  Slide {fig_num}: No images for {resp_type}")
        return slide

    images = sorted(img_dir.glob("*.png"))[:25]

    for i, img_path in enumerate(images):
        row = i // 5
        col = i % 5
        x = GRID_LEFT + col * H_STEP
        y = GRID_TOP + row * V_STEP
        slide.shapes.add_picture(str(img_path), Emu(x), Emu(y),
                                 Emu(CELL_SIZE), Emu(CELL_SIZE))

    cbar_path = FIGURES_DIR / f'Fig_{fig_num}' / f'{resp_type}_colorbar_600dpi.png'
    if cbar_path.exists():
        cbar_img = Image.open(cbar_path)
        cbar_w = int(cbar_img.width / 600 * EMU_PER_INCH)
        cbar_h = int(cbar_img.height / 600 * EMU_PER_INCH)
        grid_bottom = GRID_TOP + 4 * V_STEP + CELL_SIZE
        grid_height = grid_bottom - GRID_TOP
        cbar_x = SLIDE_WIDTH - RIGHT_MARGIN - cbar_w
        cbar_y = GRID_TOP + (grid_height - cbar_h) // 2
        slide.shapes.add_picture(str(cbar_path), Emu(cbar_x), Emu(cbar_y),
                                 Emu(cbar_w), Emu(cbar_h))

    print(f"  Slide {fig_num}: {len(images)} {resp_type} ({source}) + colorbar")
    return slide


# ============================================================================
# MAIN BUILD
# ============================================================================

def build():
    """Build the complete axisless PowerPoint."""
    with open(LAYOUT_JSON) as f:
        layout_data = json.load(f)
    layouts = layout_data['slides']

    prs = Presentation()
    prs.slide_width = Emu(SLIDE_WIDTH)
    prs.slide_height = Emu(SLIDE_HEIGHT)

    print("=" * 60)
    print("Building Axisless PowerPoint")
    print(f"  Slide size: {SLIDE_WIDTH/EMU_PER_INCH:.2f}\" x {SLIDE_HEIGHT/EMU_PER_INCH:.2f}\"")
    print("=" * 60)

    # Slide 1: Pipeline schematic (original — no axisless version)
    slide1 = _add_slide(prs, "Figure 1: Pipeline")
    fig1 = FIGURES_DIR / 'Fig_1' / 'Fig_1.png'
    if fig1.exists() and '1' in layouts:
        p = layouts['1']['Fig_1']
        _place_image(slide1, fig1, p['x'], p['y'], p['w'], p['h'])
        print("  Slide 1: Fig_1 (original)")

    # Slide 2: Figure 2 (extracted positions, row heights enforced)
    build_slide_2(prs)

    # Slide 3: Figure 3 (extracted positions, row heights enforced)
    build_slide_3(prs)

    # Slides 4-5: 5x5 grids
    build_5x5_slide(prs, '4', 'O2', "Figure 4: O2 Surface Grid")
    build_5x5_slide(prs, '5', 'Contractility', "Figure 5: Contractility Surface Grid")

    # Slides 6-8: Prediction figures (positions from slide_layout.json)
    for fig_num, target in [('6', 'Arrhythmia'), ('7', 'Heart Damage'), ('8', 'Concern')]:
        build_panel_slide(prs, fig_num,
                          ['a', 'b', 'c', 'd', 'e', 'f', 'g', 'h'],
                          layouts.get(fig_num, {}),
                          f"Figure {fig_num}: {target}")

    # Slide 9: S1 — use extracted positions
    s1_layout = layouts.get('9', {})
    build_panel_slide(prs, 'S1', ['a', 'b', 'c'], s1_layout, "Figure S1")

    # Slides 10-12: S2-S4 — auto-fit since no extracted positions
    for fig_id, panels, title in [
        ('S2', ['a', 'b'], 'Figure S2'),
        ('S3', ['a'], 'Figure S3'),
        ('S4', ['a', 'b'], 'Figure S4'),
    ]:
        slide = _add_slide(prs, title)
        placed = 0
        y_pos = 1.0
        for letter in panels:
            img_path, _ = _find_image(fig_id, letter)
            if img_path:
                img = Image.open(img_path)
                max_w = 6.0
                aspect = img.height / img.width
                w = min(max_w, 6.0)
                h = w * aspect
                x = (SLIDE_WIDTH / EMU_PER_INCH - w) / 2
                _place_image(slide, img_path, x, y_pos, w, h)
                y_pos += h + 0.3
                placed += 1
        print(f"  Slide {fig_id}: {placed} panels")

    prs.save(str(OUTPUT_PPTX))
    print(f"\n{'=' * 60}")
    print(f"Saved: {OUTPUT_PPTX}")
    print(f"{'=' * 60}")


if __name__ == '__main__':
    build()
