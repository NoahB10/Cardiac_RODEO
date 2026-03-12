"""
Build PowerPoint slides with 5x5 grid of individual images.
Uses python-pptx for image placement.

This script is located in Output/PowerPoint_Figures/ for figure management.
"""

from pptx import Presentation
from pptx.util import Inches, Emu, Pt
from pathlib import Path
import os

# Paths - this script is in Output/PowerPoint_Figures/
SCRIPT_DIR = Path(__file__).parent.resolve()
PROJECT_ROOT = SCRIPT_DIR.parent.parent  # Go up from PowerPoint_Figures -> Output -> Project Root

PPTX_PATH = SCRIPT_DIR / "Cardiac_RODEO_Tracked.pptx"
O2_DIR = SCRIPT_DIR / "Fig_4" / "O2_5x5_Individual"
CON_DIR = SCRIPT_DIR / "Fig_5" / "Contractility_5x5_Individual"
O2_CBAR = SCRIPT_DIR / "Fig_4" / "O2_colorbar_600dpi.png"
CON_CBAR = SCRIPT_DIR / "Fig_5" / "Contractility_colorbar_600dpi.png"


def build_slides():
    """Build the 5x5 grid slides in PowerPoint."""
    # Load presentation
    prs = Presentation(str(PPTX_PATH))

    # Slide dimensions (from presentation.xml: cx="6483350" cy="7745413")
    SLIDE_WIDTH = 6483350  # EMU
    SLIDE_HEIGHT = 7745413  # EMU

    # Grid positions extracted from manually-corrected layout in Tracked2.pptx
    # (averaged from first 2 rows of slide 4, which were positioned correctly)
    GRID_LEFT = 416424     # ~0.4554" — first column left edge
    GRID_TOP = 533120      # ~0.5833" — first row top edge
    H_STEP = 981300        # ~1.0730" — horizontal step between column left edges
    V_STEP = 979431        # ~1.0710" — vertical step between row top edges
    CELL_SIZE = int(1.22 * 914400)  # 1.22" square cells (transparent PNGs overlap slightly)

    # Colorbar placement
    RIGHT_MARGIN = 137160  # ~0.15" right margin for colorbar positioning

    print(f"Slide: {SLIDE_WIDTH/914400:.2f}\" x {SLIDE_HEIGHT/914400:.2f}\"")
    print(f"Cell size: {CELL_SIZE/914400:.2f}\" x {CELL_SIZE/914400:.2f}\" (square)")
    print(f"Grid origin: ({GRID_LEFT/914400:.4f}\", {GRID_TOP/914400:.4f}\")")
    print(f"H step: {H_STEP/914400:.4f}\" | V step: {V_STEP/914400:.4f}\"")

    def get_sorted_images(directory):
        """Get images sorted by their index prefix."""
        images = sorted(directory.glob("*.png"))
        return images

    def build_grid_slide(slide, image_dir, colorbar_path, title_text):
        """Build a slide with 5x5 grid of images and colorbar."""

        # Clear ALL existing shapes except the title text box
        # This removes pictures, groups (which may contain pictures from prior runs),
        # and any other non-title shapes to prevent duplicates
        shapes_to_remove = []
        for shape in slide.shapes:
            # Keep the title shape
            if hasattr(shape, "text") and title_text in shape.text:
                continue
            # Remove pictures, groups, and any other non-title shapes
            if shape.shape_type in (13, 6):  # 13=Picture, 6=Group
                shapes_to_remove.append(shape)
            elif shape.shape_type == 13:  # Picture outside group
                shapes_to_remove.append(shape)

        for shape in shapes_to_remove:
            sp = shape._element
            sp.getparent().remove(sp)

        if shapes_to_remove:
            print(f"  Removed {len(shapes_to_remove)} existing shapes")

        # Get sorted images
        images = get_sorted_images(image_dir)

        if len(images) < 25:
            print(f"Warning: Only found {len(images)} images in {image_dir}")
            return

        # Add images in 5x5 grid — 1.22" square, evenly spaced
        for i, img_path in enumerate(images[:25]):
            row = i // 5
            col = i % 5

            x = GRID_LEFT + col * H_STEP
            y = GRID_TOP + row * V_STEP

            pic = slide.shapes.add_picture(
                str(img_path),
                Emu(x), Emu(y),
                Emu(CELL_SIZE), Emu(CELL_SIZE)
            )

        # Add colorbar on the right - use actual image dimensions (no stretching)
        from PIL import Image
        with Image.open(colorbar_path) as img:
            img_width, img_height = img.size
        # At 600 DPI, convert pixels to EMU (914400 EMU = 1 inch, 600 px = 1 inch)
        cbar_width_emu = int(img_width / 600 * 914400)
        cbar_height_emu = int(img_height / 600 * 914400)

        # Place colorbar at right edge, centered vertically on grid
        grid_bottom = GRID_TOP + 4 * V_STEP + CELL_SIZE
        grid_height = grid_bottom - GRID_TOP
        cbar_x = SLIDE_WIDTH - RIGHT_MARGIN - cbar_width_emu
        cbar_y = GRID_TOP + (grid_height - cbar_height_emu) // 2

        slide.shapes.add_picture(
            str(colorbar_path),
            Emu(cbar_x), Emu(cbar_y),
            Emu(cbar_width_emu), Emu(cbar_height_emu)
        )

        print(f"Added 25 images + colorbar to slide")

    # Process slides 4 and 5 (0-indexed: 3 and 4)
    print("\n" + "="*60)
    print("Building O2 slide (Figure 4)...")
    print("="*60)
    slide4 = prs.slides[3]  # Figure 4: O2
    build_grid_slide(slide4, O2_DIR, O2_CBAR, "Figure 4")

    print("\n" + "="*60)
    print("Building Contractility slide (Figure 5)...")
    print("="*60)
    slide5 = prs.slides[4]  # Figure 5: Contractility
    build_grid_slide(slide5, CON_DIR, CON_CBAR, "Figure 5")

    # Save
    output_path = PPTX_PATH
    prs.save(str(output_path))
    print(f"\nSaved: {output_path}")
    print("\nDone! Open PowerPoint to see the high-quality 5x5 grids.")


if __name__ == '__main__':
    build_slides()
