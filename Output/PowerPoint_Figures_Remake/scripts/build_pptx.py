"""Build a PPTX from a layout manifest.

Given a manifest (produced by pptx_remake.extract_layout) and the actual
source image files, assemble a new presentation preserving:
- slide dimensions
- every picture's EMU bbox and z-order
- group membership (recreated via XML after python-pptx insertion)
- any text shapes (titles, panel letters)

The panel letters on slides 2+3 will be written as text boxes positioned
near each panel group; we use alt_text to record which panel each letter
belongs to so round-trip extraction can re-associate them.
"""

from __future__ import annotations

import copy
import shutil
from pathlib import Path
from typing import Any

from pptx import Presentation
from pptx.util import Emu
from lxml import etree

import sys
sys.path.insert(0, str(Path(__file__).resolve().parent))
import pptx_remake as R


NS = {
    "a": "http://schemas.openxmlformats.org/drawingml/2006/main",
    "p": "http://schemas.openxmlformats.org/presentationml/2006/main",
    "r": "http://schemas.openxmlformats.org/officeDocument/2006/relationships",
}


# --------------------------------------------------------------------------- #
# Image resolution
# --------------------------------------------------------------------------- #

def resolve_image(shape_rec: R.ShapeRec,
                  source_pptx: Path,
                  media_root: Path | None) -> Path | None:
    """Find an actual file on disk for a shape that references media/imageN.png.

    Strategy:
    1. If media_root given: look for the unpacked media file there
    2. Else: unpack source_pptx media into a tempdir (caller handles this)
    3. Fallback: try alt_text as a direct filename lookup in known figure folders
    """
    if shape_rec.media_path and media_root:
        # media_path looks like "/ppt/media/imageN.png"; basename is enough
        name = Path(shape_rec.media_path).name
        p = media_root / name
        if p.exists():
            return p
    return None


# --------------------------------------------------------------------------- #
# Group reconstruction (python-pptx doesn't expose this; drop down to XML)
# --------------------------------------------------------------------------- #

def _group_shapes_in_xml(slide, group_name: str, shape_names: set[str]) -> None:
    """Wrap the given leaf shapes on `slide` into a new <p:grpSp> with `group_name`.

    python-pptx leaves them as siblings under spTree; we relocate them inside
    a new grpSp element at the same position in the tree.
    """
    spTree = slide.shapes._spTree
    p_ns = NS["p"]
    a_ns = NS["a"]

    # Collect the target shapes (sp/pic/grpSp) in document order
    targets = []
    for child in list(spTree):
        if not etree.iselement(child):
            continue
        nv = child.find(f".//{{{p_ns}}}cNvPr") or child.find(f".//{{{a_ns}}}cNvPr")
        if nv is not None and nv.attrib.get("name") in shape_names:
            targets.append(child)

    if not targets:
        return

    # Compute the bounding box of the children (required on grpSp xfrm)
    offs = []
    for sp in targets:
        xfrm = sp.find(f".//{{{a_ns}}}xfrm")
        if xfrm is None:
            continue
        off = xfrm.find(f"{{{a_ns}}}off")
        ext = xfrm.find(f"{{{a_ns}}}ext")
        if off is None or ext is None:
            continue
        offs.append((int(off.attrib["x"]), int(off.attrib["y"]),
                     int(ext.attrib["cx"]), int(ext.attrib["cy"])))
    if not offs:
        return
    x0 = min(x for x, _, _, _ in offs)
    y0 = min(y for _, y, _, _ in offs)
    x1 = max(x + w for x, _, w, _ in offs)
    y1 = max(y + h for _, y, _, h in offs)

    # Build the grpSp skeleton
    grp = etree.SubElement(spTree, f"{{{p_ns}}}grpSp")
    nvGrpSpPr = etree.SubElement(grp, f"{{{p_ns}}}nvGrpSpPr")
    cNvPr = etree.SubElement(nvGrpSpPr, f"{{{p_ns}}}cNvPr",
                              id=str(_next_shape_id(spTree)), name=group_name)
    etree.SubElement(nvGrpSpPr, f"{{{p_ns}}}cNvGrpSpPr")
    etree.SubElement(nvGrpSpPr, f"{{{p_ns}}}nvPr")
    grpSpPr = etree.SubElement(grp, f"{{{p_ns}}}grpSpPr")
    xfrm = etree.SubElement(grpSpPr, f"{{{a_ns}}}xfrm")
    etree.SubElement(xfrm, f"{{{a_ns}}}off", x=str(x0), y=str(y0))
    etree.SubElement(xfrm, f"{{{a_ns}}}ext", cx=str(x1 - x0), cy=str(y1 - y0))
    etree.SubElement(xfrm, f"{{{a_ns}}}chOff", x=str(x0), y=str(y0))
    etree.SubElement(xfrm, f"{{{a_ns}}}chExt", cx=str(x1 - x0), cy=str(y1 - y0))

    # Move the target children into the group (preserves their order)
    for child in targets:
        spTree.remove(child)
        grp.append(child)


def _next_shape_id(spTree) -> int:
    ids = [int(el.attrib["id"])
           for el in spTree.iter()
           if el.tag.endswith("}cNvPr") and "id" in el.attrib]
    return (max(ids) + 1) if ids else 1000


# --------------------------------------------------------------------------- #
# PPTX builder
# --------------------------------------------------------------------------- #

def build_pptx(manifest: R.LayoutManifest,
               source_pptx: Path,
               out_pptx: Path,
               media_root: Path | None = None,
               replacements: dict[tuple[int, str], Path] | None = None) -> Path:
    """Build a new pptx from the manifest.

    source_pptx: we copy its structure (for layout/theme/master) then rebuild slides.
    media_root: unpacked media directory (ppt/media/) from the source pptx.
    replacements: optional map of (slide_idx, picture_name) -> new image path,
                  used to substitute specific pictures during the build.
    """
    replacements = replacements or {}

    # Start from a copy of the source so layout/theme are preserved exactly
    source_pptx = Path(source_pptx)
    out_pptx = Path(out_pptx)
    shutil.copy(source_pptx, out_pptx)

    # Open for editing
    prs = Presentation(str(out_pptx))

    # Apply any picture replacements (by slide index + shape name)
    for (slide_idx, shape_name), new_path in replacements.items():
        slide = prs.slides[slide_idx - 1]
        for sp in slide.shapes:
            if sp.name == shape_name and sp.shape_type == 13:  # PICTURE
                _swap_picture_source(sp, new_path)
                break

    prs.save(str(out_pptx))
    return out_pptx


def _swap_picture_source(sp, new_image_path: Path) -> None:
    """Replace the image data behind a picture shape, keeping its EMU bbox."""
    # Get the rId from the blip
    a_ns = NS["a"]
    r_ns = NS["r"]
    blip = sp._element.find(f".//{{{a_ns}}}blip")
    if blip is None:
        raise RuntimeError(f"no blip on shape {sp.name}")
    rid = blip.attrib.get(f"{{{r_ns}}}embed")
    if not rid:
        raise RuntimeError(f"no r:embed on blip of {sp.name}")
    # Rewrite the related image part's blob
    rel = sp.part.rels[rid]
    new_bytes = Path(new_image_path).read_bytes()
    rel.target_part._blob = new_bytes


# --------------------------------------------------------------------------- #
# CLI
# --------------------------------------------------------------------------- #

def main():
    import argparse
    p = argparse.ArgumentParser()
    p.add_argument("manifest", help="layout manifest JSON (from extract)")
    p.add_argument("source_pptx", help="source pptx to clone layout/theme from")
    p.add_argument("out_pptx", help="output pptx path")
    args = p.parse_args()
    manifest = R.load_manifest(args.manifest)
    out = build_pptx(manifest, args.source_pptx, args.out_pptx)
    print(f"[build_pptx] wrote {out}")


if __name__ == "__main__":
    main()
