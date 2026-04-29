"""Add auto-placed, row-aligned panel letters to a tracked pptx.

For each slide, identify the panels (groups + standalone pictures above a
size threshold), cluster them into rows by y-coordinate, assign letters
a, b, c, ... in reading order, and add a text-box shape for each letter
just above-left of its panel. All letters in the same row share an
identical y; letters in the same column share an identical x.

Slides to skip:
- Slide 1 (title figure)
- Slides 4, 5 (5×5 surface grids — each slide is one figure, not panels)

For slides 2+3 the existing groups define the panels. For 6-12 every
picture larger than MIN_PANEL_IN is a panel.
"""

from __future__ import annotations

from pathlib import Path
from typing import Any

from pptx import Presentation
from pptx.util import Emu, Pt, Inches
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from lxml import etree

EMU_PER_INCH = 914400
SKIP_SLIDES = {1, 4, 5}
MIN_PANEL_IN = 0.5          # pictures smaller than this ignored as panels
ROW_CLUSTER_IN = 0.4        # panels whose tops are within this inch are same row
COL_CLUSTER_IN = 0.4        # panels whose lefts are within this inch are same col
LETTER_OFFSET_X_IN = 0.05   # letter sits this far left of panel's left edge
LETTER_OFFSET_Y_IN = 0.15   # letter sits this far above panel's top edge
LETTER_BOX_W_IN = 0.30
LETTER_BOX_H_IN = 0.18
LETTER_FONT_SIZE_PT = 12
LETTER_FONT_NAME = "Helvetica"
LETTER_FONT_BOLD = False
LETTER_COLOR = RGBColor(0x00, 0x00, 0x00)

ALPHABET = "abcdefghijklmnopqrstuvwxyz"
NS = {
    "a": "http://schemas.openxmlformats.org/drawingml/2006/main",
    "p": "http://schemas.openxmlformats.org/presentationml/2006/main",
}


# --------------------------------------------------------------------------- #
# Panel discovery
# --------------------------------------------------------------------------- #

def _shape_bbox(shape) -> tuple[int, int, int, int]:
    return (int(shape.left or 0), int(shape.top or 0),
            int(shape.width or 0), int(shape.height or 0))


def _find_panels(slide) -> list[tuple[str, int, int, int, int]]:
    """Return a list of (panel_id, left, top, width, height) EMU tuples.

    panel_id is either the group's name or the picture's name.
    Groups are treated as single panels; standalone pictures are panels too
    iff larger than MIN_PANEL_IN on both dimensions.
    """
    panels: list[tuple[str, int, int, int, int]] = []
    min_emu = int(MIN_PANEL_IN * EMU_PER_INCH)
    seen_in_groups: set[str] = set()

    for shape in slide.shapes:
        # Groups -> one panel with group's overall bbox
        if shape.shape_type == 6:  # MSO_SHAPE_TYPE.GROUP
            L, T, W, H = _shape_bbox(shape)
            panels.append((shape.name, L, T, W, H))
            # mark children
            for child in shape.shapes:
                seen_in_groups.add(child.name)

    for shape in slide.shapes:
        if shape.shape_type == 13:  # PICTURE
            if shape.name in seen_in_groups:
                continue
            L, T, W, H = _shape_bbox(shape)
            if W >= min_emu and H >= min_emu:
                panels.append((shape.name, L, T, W, H))

    return panels


# --------------------------------------------------------------------------- #
# Row / column clustering and alignment
# --------------------------------------------------------------------------- #

def _cluster(values: list[int], tol_emu: int) -> dict[int, int]:
    """Cluster near-equal values; return value -> cluster_representative."""
    sorted_vals = sorted(set(values))
    clusters: list[list[int]] = []
    for v in sorted_vals:
        if clusters and v - clusters[-1][-1] <= tol_emu:
            clusters[-1].append(v)
        else:
            clusters.append([v])
    mapping: dict[int, int] = {}
    for c in clusters:
        rep = int(sum(c) / len(c))
        for v in c:
            mapping[v] = rep
    return mapping


def _align_panels(panels) -> list[tuple]:
    """Snap panels so same-row tops and same-column lefts match."""
    if not panels:
        return panels
    tol_row = int(ROW_CLUSTER_IN * EMU_PER_INCH)
    tol_col = int(COL_CLUSTER_IN * EMU_PER_INCH)
    tops = _cluster([p[2] for p in panels], tol_row)
    lefts = _cluster([p[1] for p in panels], tol_col)
    snapped = []
    for name, L, T, W, H in panels:
        snapped.append((name, lefts[L], tops[T], W, H))
    return snapped


def _reading_order(panels) -> list[tuple]:
    """Sort panels top-left to bottom-right (after snapping)."""
    snapped = _align_panels(panels)
    return sorted(snapped, key=lambda p: (p[2], p[1]))


# --------------------------------------------------------------------------- #
# Letter placement
# --------------------------------------------------------------------------- #

def _add_letter_textbox(slide, letter: str, x_emu: int, y_emu: int,
                         panel_name: str) -> None:
    w = int(LETTER_BOX_W_IN * EMU_PER_INCH)
    h = int(LETTER_BOX_H_IN * EMU_PER_INCH)
    tb = slide.shapes.add_textbox(x_emu, y_emu, w, h)
    tb.name = f"PanelLetter_{letter}"
    # Tag the letter with the panel it belongs to, for later regrouping
    tb._element.nvSpPr.cNvPr.attrib["descr"] = f"panel={panel_name}"
    tf = tb.text_frame
    tf.margin_left = 0
    tf.margin_right = 0
    tf.margin_top = 0
    tf.margin_bottom = 0
    p = tf.paragraphs[0]
    p.alignment = PP_ALIGN.LEFT
    run = p.add_run()
    run.text = letter
    run.font.name = LETTER_FONT_NAME
    run.font.size = Pt(LETTER_FONT_SIZE_PT)
    run.font.bold = LETTER_FONT_BOLD
    run.font.color.rgb = LETTER_COLOR


def place_letters_on_slide(slide, slide_index: int) -> int:
    """Add letters to one slide. Returns count added."""
    panels = _find_panels(slide)
    if not panels:
        return 0
    ordered = _reading_order(panels)
    count = 0
    for idx, (name, L, T, W, H) in enumerate(ordered):
        if idx >= len(ALPHABET):
            break
        letter = ALPHABET[idx]
        x = L - int(LETTER_OFFSET_X_IN * EMU_PER_INCH)
        y = T - int(LETTER_OFFSET_Y_IN * EMU_PER_INCH)
        if x < 0:
            x = 0
        if y < 0:
            y = 0
        _add_letter_textbox(slide, letter, x, y, name)
        count += 1
    return count


def place_letters(pptx_path: str | Path, out_path: str | Path) -> Path:
    prs = Presentation(str(pptx_path))
    total = 0
    for i, slide in enumerate(prs.slides, start=1):
        if i in SKIP_SLIDES:
            continue
        n = place_letters_on_slide(slide, i)
        total += n
        print(f"  Slide {i}: +{n} letter(s)")
    out_path = Path(out_path)
    prs.save(str(out_path))
    print(f"[place_letters] {total} letters -> {out_path}")
    return out_path


# --------------------------------------------------------------------------- #
# CLI
# --------------------------------------------------------------------------- #

def main():
    import argparse
    p = argparse.ArgumentParser()
    p.add_argument("in_pptx")
    p.add_argument("out_pptx")
    args = p.parse_args()
    place_letters(args.in_pptx, args.out_pptx)


if __name__ == "__main__":
    main()
