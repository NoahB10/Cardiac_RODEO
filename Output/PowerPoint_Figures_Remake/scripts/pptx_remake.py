"""Core module for the PowerPoint_Figures_Remake tracked system.

Handles the three pipeline phases:
    1. extract_layout(pptx_path)   — pptx -> layout manifest (JSON)
    2. render_axis_overlay(row)    — Axis_Scaling_Reference row -> transparent axis PNG
    3. build_pptx(layout, out)     — layout + panel images -> new pptx

Design:
- Layout is a plain dict / JSON (slides -> shapes) so it round-trips by diff.
- Position data is kept in EMU internally; inch columns are convenience only.
- Two-pic groups are treated as {axisless, overlay}; three-pic adds {legend}.
- Panel letters are separate text shapes aligned to the panel group.
"""

from __future__ import annotations

import json
import re
from dataclasses import dataclass, field, asdict
from pathlib import Path
from typing import Any

from pptx import Presentation
from pptx.util import Emu, Inches
from pptx.enum.shapes import MSO_SHAPE_TYPE

EMU_PER_INCH = 914400


# --------------------------------------------------------------------------- #
# Data model
# --------------------------------------------------------------------------- #

@dataclass
class ShapeRec:
    name: str
    kind: str               # "picture" | "text" | "group" | "other"
    left_emu: int
    top_emu: int
    width_emu: int
    height_emu: int
    z_order: int            # index within spTree (0-based)
    group_name: str | None = None
    alt_text: str | None = None
    rid: str | None = None
    media_path: str | None = None
    text: str | None = None

    def bbox_in(self) -> tuple[float, float, float, float]:
        return (self.left_emu / EMU_PER_INCH,
                self.top_emu / EMU_PER_INCH,
                self.width_emu / EMU_PER_INCH,
                self.height_emu / EMU_PER_INCH)


@dataclass
class SlideRec:
    index: int              # 1-based slide number
    width_emu: int
    height_emu: int
    shapes: list[ShapeRec] = field(default_factory=list)


@dataclass
class LayoutManifest:
    pptx_source: str
    slide_width_emu: int
    slide_height_emu: int
    slides: list[SlideRec] = field(default_factory=list)

    def to_dict(self) -> dict:
        return {
            "pptx_source": self.pptx_source,
            "slide_width_emu": self.slide_width_emu,
            "slide_height_emu": self.slide_height_emu,
            "slides": [
                {
                    "index": s.index,
                    "shapes": [asdict(sh) for sh in s.shapes],
                }
                for s in self.slides
            ],
        }

    @classmethod
    def from_dict(cls, d: dict) -> "LayoutManifest":
        m = cls(
            pptx_source=d["pptx_source"],
            slide_width_emu=d["slide_width_emu"],
            slide_height_emu=d["slide_height_emu"],
        )
        for s in d["slides"]:
            rec = SlideRec(
                index=s["index"],
                width_emu=d["slide_width_emu"],
                height_emu=d["slide_height_emu"],
                shapes=[ShapeRec(**sh) for sh in s["shapes"]],
            )
            m.slides.append(rec)
        return m


# --------------------------------------------------------------------------- #
# Phase 1: extract_layout
# --------------------------------------------------------------------------- #

def _walk_shapes(shapes, group_name=None, z=[0]):
    for sp in shapes:
        if sp.shape_type == MSO_SHAPE_TYPE.GROUP:
            name = sp.name
            yield ("group", sp, group_name, z[0])
            z[0] += 1
            yield from _walk_shapes(sp.shapes, name, z)
        else:
            yield ("leaf", sp, group_name, z[0])
            z[0] += 1


def _pic_media_path(sp) -> str | None:
    try:
        # sp.image.blob exists; map via sp.part.related_parts for target
        for rid, part in sp.part.rels.items():
            if part.reltype.endswith("/image") and part.target_ref in str(sp._element.xml):
                return part.target_ref
    except Exception:
        pass
    # Fallback: via r:embed on blipFill
    try:
        ns = {"a": "http://schemas.openxmlformats.org/drawingml/2006/main",
              "r": "http://schemas.openxmlformats.org/officeDocument/2006/relationships"}
        blip = sp._element.find(".//a:blip", ns)
        if blip is not None:
            rid = blip.attrib.get("{http://schemas.openxmlformats.org/officeDocument/2006/relationships}embed")
            if rid:
                part = sp.part.rels[rid].target_part
                return part.partname
    except Exception:
        pass
    return None


def _pic_rid(sp) -> str | None:
    try:
        ns = {"a": "http://schemas.openxmlformats.org/drawingml/2006/main"}
        blip = sp._element.find(".//a:blip", ns)
        if blip is not None:
            return blip.attrib.get(
                "{http://schemas.openxmlformats.org/officeDocument/2006/relationships}embed"
            )
    except Exception:
        pass
    return None


def _shape_to_rec(sp, group_name: str | None, z: int) -> ShapeRec:
    kind = "other"
    alt = None
    rid = None
    media_path = None
    text = None

    if sp.shape_type == MSO_SHAPE_TYPE.GROUP:
        kind = "group"
    elif sp.shape_type == MSO_SHAPE_TYPE.PICTURE:
        kind = "picture"
        try:
            alt = sp._element.nvPicPr.cNvPr.attrib.get("descr", "") or None
        except Exception:
            alt = None
        rid = _pic_rid(sp)
        media_path = _pic_media_path(sp)
    elif getattr(sp, "has_text_frame", False):
        kind = "text"
        try:
            text = sp.text_frame.text
        except Exception:
            text = None

    return ShapeRec(
        name=sp.name,
        kind=kind,
        left_emu=int(sp.left or 0),
        top_emu=int(sp.top or 0),
        width_emu=int(sp.width or 0),
        height_emu=int(sp.height or 0),
        z_order=z,
        group_name=group_name,
        alt_text=alt,
        rid=rid,
        media_path=media_path,
        text=text,
    )


def extract_layout(pptx_path: str | Path) -> LayoutManifest:
    """Parse pptx into a layout manifest. Pure read-only."""
    pptx_path = Path(pptx_path)
    prs = Presentation(str(pptx_path))
    manifest = LayoutManifest(
        pptx_source=str(pptx_path),
        slide_width_emu=int(prs.slide_width),
        slide_height_emu=int(prs.slide_height),
    )
    for i, slide in enumerate(prs.slides, start=1):
        rec = SlideRec(
            index=i,
            width_emu=int(prs.slide_width),
            height_emu=int(prs.slide_height),
        )
        z_counter = [0]
        for tag, sp, gname, _ in _walk_shapes(slide.shapes, None, z_counter):
            # skip group container itself; keep only leaves + record group membership
            if tag == "group":
                continue
            rec.shapes.append(_shape_to_rec(sp, gname, z_counter[0] - 1))
        manifest.slides.append(rec)
    return manifest


def save_manifest(manifest: LayoutManifest, out_path: str | Path) -> Path:
    out_path = Path(out_path)
    out_path.write_text(json.dumps(manifest.to_dict(), indent=2))
    return out_path


def load_manifest(path: str | Path) -> LayoutManifest:
    return LayoutManifest.from_dict(json.loads(Path(path).read_text()))


# --------------------------------------------------------------------------- #
# Phase 1b: panel identification (which groups are overlay-eligible)
# --------------------------------------------------------------------------- #

@dataclass
class PanelGroup:
    slide: int
    group_name: str
    pictures: list[ShapeRec]
    # By convention: smallest area = axisless base, largest = overlay,
    #                tiny inside bbox = legend
    def classify(self) -> dict[str, ShapeRec | None]:
        pics = sorted(self.pictures, key=lambda s: s.width_emu * s.height_emu)
        out = {"axisless": None, "overlay": None, "legend": None}
        if len(pics) == 0:
            return out
        if len(pics) == 1:
            out["axisless"] = pics[0]
        elif len(pics) == 2:
            out["axisless"], out["overlay"] = pics[0], pics[1]
        else:  # 3+
            out["legend"] = pics[0]          # smallest
            out["axisless"] = pics[1]        # middle
            out["overlay"] = pics[-1]        # largest
        return out


def find_panel_groups(manifest: LayoutManifest) -> list[PanelGroup]:
    groups: dict[tuple[int, str], PanelGroup] = {}
    for slide in manifest.slides:
        for sh in slide.shapes:
            if sh.kind != "picture" or not sh.group_name:
                continue
            key = (slide.index, sh.group_name)
            if key not in groups:
                groups[key] = PanelGroup(
                    slide=slide.index, group_name=sh.group_name, pictures=[]
                )
            groups[key].pictures.append(sh)
    return list(groups.values())


# --------------------------------------------------------------------------- #
# Phase 2: render_axis_overlay (stub for next step)
# --------------------------------------------------------------------------- #

def render_axis_overlay(axis_row: dict, out_path: str | Path) -> Path:
    """Render a transparent axis-only PNG sized to match the axisless image.
    See scripts/axis_overlay.py for the full implementation.
    """
    from axis_overlay import render_overlay
    return render_overlay(axis_row, out_path)


# --------------------------------------------------------------------------- #
# CLI
# --------------------------------------------------------------------------- #

def _cli_extract(args):
    manifest = extract_layout(args.pptx)
    out = save_manifest(manifest, args.out)
    print(f"[extract_layout] {len(manifest.slides)} slides -> {out}")
    for s in manifest.slides:
        npic = sum(1 for sh in s.shapes if sh.kind == "picture")
        ngrp = len({sh.group_name for sh in s.shapes if sh.group_name})
        print(f"  Slide {s.index}: {npic} pictures, {ngrp} groups")


def main():
    import argparse
    p = argparse.ArgumentParser()
    sub = p.add_subparsers(dest="cmd", required=True)

    e = sub.add_parser("extract", help="extract layout manifest from pptx")
    e.add_argument("pptx")
    e.add_argument("out")
    e.set_defaults(func=_cli_extract)

    args = p.parse_args()
    args.func(args)


if __name__ == "__main__":
    main()
