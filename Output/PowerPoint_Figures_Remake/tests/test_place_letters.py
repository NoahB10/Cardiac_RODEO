"""Tests for Phase 4: place_letters."""

import sys
from pathlib import Path

import pytest
from pptx import Presentation

HERE = Path(__file__).resolve().parent
SCRIPTS = HERE.parent / "scripts"
sys.path.insert(0, str(SCRIPTS))

import place_letters as PL  # noqa: E402
import pptx_remake as R      # noqa: E402

PROJECT_ROOT = HERE.parents[2]
SOURCE_PPTX = PROJECT_ROOT / "Output" / "PowerPoint_Figures" / "Cardiac_RODEO.pptx"


def _letters_per_slide(pptx_path):
    prs = Presentation(str(pptx_path))
    result = {}
    for i, slide in enumerate(prs.slides, start=1):
        letters = []
        for sh in slide.shapes:
            if sh.name.startswith("PanelLetter_"):
                letter = sh.name.replace("PanelLetter_", "")
                letters.append((letter, int(sh.left), int(sh.top)))
        result[i] = letters
    return result


def test_letters_added_to_expected_slides(tmp_path):
    out = tmp_path / "lettered.pptx"
    PL.place_letters(SOURCE_PPTX, out)
    by_slide = _letters_per_slide(out)
    # Skipped: 1, 4, 5
    assert by_slide[1] == [], "slide 1 should have no letters"
    assert by_slide[4] == [], "slide 4 should have no letters"
    assert by_slide[5] == [], "slide 5 should have no letters"


def test_letters_in_reading_order(tmp_path):
    out = tmp_path / "lettered.pptx"
    PL.place_letters(SOURCE_PPTX, out)
    by_slide = _letters_per_slide(out)
    # For every slide with letters, the alphabetic order must match sort-by(top, left)
    for slide_idx, letters in by_slide.items():
        if not letters:
            continue
        # letters as-added should be in alphabetical order
        sorted_by_pos = sorted(letters, key=lambda t: (t[2], t[1]))  # sort by (top, left)
        alpha_order = sorted(letters, key=lambda t: t[0])
        assert [l[0] for l in sorted_by_pos] == [l[0] for l in alpha_order], \
            f"Slide {slide_idx} letters not in reading order: {letters}"


def test_letters_row_aligned(tmp_path):
    """Letters in the same row must have identical y-coordinates."""
    out = tmp_path / "lettered.pptx"
    PL.place_letters(SOURCE_PPTX, out)
    by_slide = _letters_per_slide(out)
    tol_emu = int(PL.ROW_CLUSTER_IN * PL.EMU_PER_INCH)
    for slide_idx, letters in by_slide.items():
        if len(letters) < 2:
            continue
        # Group by y (within tolerance)
        from collections import defaultdict
        groups = defaultdict(list)
        for l, x, y in letters:
            # Bucket to 0.01 inch to cluster
            bucket = round(y / tol_emu) * tol_emu
            groups[bucket].append((l, x, y))
        for bucket, grp in groups.items():
            ys = [y for _, _, y in grp]
            assert max(ys) - min(ys) < 10, \
                f"Slide {slide_idx} row at y≈{bucket}: ys differ by >10 EMU: {ys}"


def test_bbox_unchanged_after_letters(tmp_path):
    """Adding letters must not move any picture."""
    out = tmp_path / "lettered.pptx"
    PL.place_letters(SOURCE_PPTX, out)
    src_manifest = R.extract_layout(SOURCE_PPTX)
    dst_manifest = R.extract_layout(out)
    src_pics = {(s.index, sh.name): sh
                for s in src_manifest.slides for sh in s.shapes
                if sh.kind == "picture"}
    dst_pics = {(s.index, sh.name): sh
                for s in dst_manifest.slides for sh in s.shapes
                if sh.kind == "picture"}
    assert set(src_pics) == set(dst_pics)
    for k, sp in src_pics.items():
        dp = dst_pics[k]
        assert (sp.left_emu, sp.top_emu, sp.width_emu, sp.height_emu) == \
               (dp.left_emu, dp.top_emu, dp.width_emu, dp.height_emu), \
               f"picture moved: {k}"


def test_letters_link_to_panels_via_alt_text(tmp_path):
    """Each letter text box has descr='panel=<name>' for later regrouping."""
    out = tmp_path / "lettered.pptx"
    PL.place_letters(SOURCE_PPTX, out)
    prs = Presentation(str(out))
    for slide in prs.slides:
        for sh in slide.shapes:
            if sh.name.startswith("PanelLetter_"):
                descr = sh._element.nvSpPr.cNvPr.attrib.get("descr", "")
                assert descr.startswith("panel="), \
                    f"letter {sh.name} missing panel link"


def test_panel_count_matches_detection(tmp_path):
    """The total letter count must equal the detected panel count (non-skipped slides, ≤26/slide)."""
    out = tmp_path / "lettered.pptx"
    PL.place_letters(SOURCE_PPTX, out)
    prs_src = Presentation(str(SOURCE_PPTX))
    expected = 0
    for i, slide in enumerate(prs_src.slides, start=1):
        if i in PL.SKIP_SLIDES:
            continue
        n = len(PL._find_panels(slide))
        expected += min(n, 26)
    by_slide = _letters_per_slide(out)
    actual = sum(len(v) for v in by_slide.values())
    assert actual == expected, \
        f"expected {expected} letters, got {actual}"
