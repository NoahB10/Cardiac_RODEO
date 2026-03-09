"""
Update Cardiac RODEO Tracked PowerPoint
Uses actual saved sizes from registry - no stretching.
Adds detailed titles and slide notes for each figure.
"""

from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN
from pathlib import Path
import pandas as pd

# Paths
PROJECT_ROOT = Path(__file__).parent
FIGURES_DIR = PROJECT_ROOT / 'Output' / 'PowerPoint_Figures'
PPTX_PATH = FIGURES_DIR / 'Cardiac_RODEO_Tracked.pptx'
REGISTRY_PATH = FIGURES_DIR / 'figure_registry.csv'

# Slide dimensions (portrait 7.09" x 8.47")
SLIDE_WIDTH = 7.09
SLIDE_HEIGHT = 8.47
MARGIN = 0.3
GAP = 0.15
TITLE_HEIGHT = 0.4

# ============================================================================
# FIGURE CONFIGURATIONS WITH TITLES AND NOTES
# ============================================================================

FIGURE_CONFIGS = {
    '1': {
        'title': 'Figure 1: Cardiac RODEO Pipeline Overview',
        'notes': """FIGURE 1: PIPELINE SCHEMATIC

PURPOSE: Provide a visual overview of the complete Cardiac RODEO workflow.

CONTENT:
- Organoid culture and drug exposure setup
- Dynamix oxygen sensor data acquisition
- Signal processing pipeline (SNR filtering, contractility extraction)
- PK-PD elimination equation fitting
- Feature extraction (14 coefficients: 7 per response type)
- Machine learning classification (Arrhythmia, Heart Damage, Concern)

STATUS: Placeholder - requires external schematic (BioRender or similar)

DATA SOURCE: N/A (schematic)
"""
    },
    '2': {
        'title': 'Figure 2: Robust Data Generation & Quality Control',
        'notes': """FIGURE 2: ROBUST GENERATION

PANEL a - SNR Distribution Analysis:
- Bar plot showing SNR distribution across all measurements
- Threshold at 0.4 (vertical line)
- Blue bars = PASS (SNR ≥ 0.4), Grey bars = FAIL (SNR < 0.4)
- Square aspect ratio (2.9" x 2.9")
- Shows data quality filtering criteria

PANEL b - External Images:
- Placeholder for experimental images
- Organoid microscopy, contractility traces, oxygen sensor setup

DATA SOURCES:
- Output/Excel_Figures/snr_analysis.xlsx
- External microscopy images (to be added)

GOAL: Demonstrate robust signal quality and filtering methodology.
"""
    },
    '3': {
        'title': 'Figure 3: PK-PD Surface Fitting & Equation Selection',
        'notes': """FIGURE 3: FITTING KINETICS

PANEL a - Vandetanib O2 Mean Heatmap:
- 2D heatmap of O2 response vs time and concentration
- Drug: Vandetanib (representative example)
- RdBu_r colormap (blue=low, red=high)

PANEL b - Vandetanib Contractility (Amp_std) Heatmap:
- 2D heatmap of contractility amplitude standard deviation
- Same drug, parallel response type

PANEL c - 3D Surface Fit:
- Placeholder for 3D PK-PD elimination surface
- Shows fitted equation overlaid on raw data
- View angle: elev=25, azim=-158

PANEL d - R² Equation Comparison:
- Horizontal bar chart comparing 12 equations
- Blue = Contractility, Pink = O2
- Shows which equation best fits the data
- PKPD Elimination (Eq11) selected for final model

PANEL e - Random Forest vs Equations:
- Comparison of RF feature-based prediction vs equation coefficients
- Demonstrates that fitted coefficients capture biological signal

DATA SOURCES:
- Cleaned_Data/Heatmaps/Vandetanib (G11)/*.csv
- Output/Excel_Figures/r2_equation_comparison.xlsx
- EQN_Coefficients/all_equations_coefficients.xlsx

GOAL: Justify equation selection and demonstrate fitting quality.
"""
    },
    '4': {
        'title': 'Figure 4: O2 Response Surfaces (All 25 Drugs)',
        'notes': """FIGURE 4: O2 3D SURFACE GRID

CONTENT:
- 5×5 grid of 3D PK-PD elimination surfaces
- All 25 drugs in the study
- Response type: O2 (oxygen consumption)

AXES:
- X: Time (0-96 hours)
- Y: Dose ratio (C0/Cmax, 0-2)
- Z: O2 response (normalized)

COLORMAP: viridis (consistent across all surfaces)
VIEW ANGLE: elev=25, azim=-158

DATA SOURCE:
- EQN_Coefficients/all_equations_coefficients.xlsx (pkpd_elimination sheet)

GOAL: Visual overview of drug-specific O2 response kinetics.
"""
    },
    '5': {
        'title': 'Figure 5: Contractility Response Surfaces (All 25 Drugs)',
        'notes': """FIGURE 5: CONTRACTILITY 3D SURFACE GRID

CONTENT:
- 5×5 grid of 3D PK-PD elimination surfaces
- All 25 drugs in the study
- Response type: Contractility (Amp_std)

AXES:
- X: Time (0-96 hours)
- Y: Dose ratio (C0/Cmax, 0-2)
- Z: Contractility response (normalized)

COLORMAP: viridis (consistent across all surfaces)
VIEW ANGLE: elev=25, azim=-158

DATA SOURCE:
- EQN_Coefficients/all_equations_coefficients.xlsx (pkpd_elimination sheet)

GOAL: Visual overview of drug-specific contractility kinetics.
"""
    },
    '6': {
        'title': 'Figure 6: Arrhythmia Prediction Model Performance',
        'notes': """FIGURE 6: ARRHYTHMIA PREDICTION

PANEL a - ROC Curve:
- Mean ROC with ±1 std shading (10-fold LOOCV)
- Blue curve, AUC displayed in legend
- Random classifier line (dashed)

PANEL b - Confusion Matrix:
- 2×2 matrix with counts and percentages
- Blues colormap
- Labels: Actual Neg/Pos vs Predicted Neg/Pos

PANEL c - Performance Metrics Bar:
- Grouped bars: Accuracy, AUC, F1, MCC
- Error bars from cross-validation
- Colors: blue, pink, orange, beige

PANEL d - Threshold Analysis:
- Horizontal scatter: drugs on Y-axis, probability on X-axis
- Blue = positive (arrhythmia), Grey = negative
- Threshold line at optimal cutoff (35%)

PANEL e - Cumulative Feature Importance:
- Shows prediction improvement as features added
- Blue = pass threshold, Grey = fail threshold
- Demonstrates which coefficients matter most

PANEL f - SHAP Aligned Pairs:
- Horizontal lines from 0: positive SHAP right, negative left
- Color = actual class (blue=arrhythmia, grey=no arrhythmia)
- Top 5 features by |mean SHAP|

PANEL g - MoLFormer ROC Comparison:
- Compare Organoid model vs CNN models
- CNN (DIQT Transfer), CNN (5-fold on 25)

PANEL h - MoLFormer Metrics Comparison:
- Grouped bars: Accuracy, Sensitivity, Specificity
- All three models compared

DATA SOURCES:
- Output/ROC_Data/roc_curves_all_models.xlsx
- Output/Confusion_Matrices/confusion_matrices_all.xlsx
- Output/SHAP_Data/shap_arrhythmia_values.csv
- Output/MoLFormer_Comparison/comparison_metrics_all.csv

GOAL: Comprehensive evaluation of arrhythmia classification.
"""
    },
    '7': {
        'title': 'Figure 7: Heart Damage Prediction Model Performance',
        'notes': """FIGURE 7: HEART DAMAGE PREDICTION

PANEL a - ROC Curve:
- Mean ROC with ±1 std shading (10-fold LOOCV)
- Blue curve, AUC displayed in legend

PANEL b - Confusion Matrix:
- 2×2 matrix with counts and percentages

PANEL c - Performance Metrics Bar:
- Grouped bars: Accuracy, AUC, F1, MCC

PANEL d - Threshold Analysis:
- Horizontal scatter with threshold at 5%
- Blue = heart damage, Grey = no damage

PANEL e - Cumulative Feature Importance:
- Feature addition analysis

PANEL f - SHAP Aligned Pairs:
- Top 5 features, blue=damage, grey=no damage

PANEL g - ADMET ROC Comparison:
- Compare Organoid vs ADMET-AI vs SwissADME
- DICTrank-based predictions

PANEL h - ADMET Metrics Comparison:
- Accuracy comparison across all models
- Pink = ADMET tools, Blue = Organoid models

DATA SOURCES:
- Output/ROC_Data/roc_curves_all_models.xlsx
- Output/Confusion_Matrices/confusion_matrices_all.xlsx
- Output/SHAP_Data/shap_heart_damage_values.csv
- Output/ADMET_Comparison/final_comparison_dictrank_vs_organoid.csv

GOAL: Heart damage classification with ADMET benchmark comparison.
"""
    },
    '8': {
        'title': 'Figure 8: Concern Level Prediction (Binary)',
        'notes': """FIGURE 8: CONCERN BINARY PREDICTION

PANEL a - ROC Curve:
- Mean ROC with ±1 std shading
- Binary: Most Concern vs Less/No Concern

PANEL b - Confusion Matrix:
- 2×2 matrix with counts and percentages

PANEL c - Performance Metrics Bar:
- Grouped bars: Accuracy, AUC, F1, MCC

PANEL d - Threshold Analysis:
- Horizontal scatter with threshold at 40%
- Blue = most concern, Grey = less/no concern

PANEL e - Cumulative Feature Importance:
- Feature addition analysis

PANEL f - SHAP Aligned Pairs:
- Top 5 features by importance

NOTE: No external model comparison for Concern
(MoLFormer and ADMET not applicable to this endpoint)

DATA SOURCES:
- Output/ROC_Data/roc_curves_all_models.xlsx
- Output/Confusion_Matrices/concern_binary_confusion_matrix.csv
- Output/SHAP_Data/shap_concern_binary_values.csv

GOAL: Clinical concern level classification from organoid data.
"""
    },
    'S1': {
        'title': 'Figure S1: Supplementary Vandetanib Heatmaps',
        'notes': """FIGURE S1: SUPPLEMENTARY HEATMAPS

Additional response metrics for Vandetanib (representative drug).

PANEL a - O2_std:
- Standard deviation of O2 signal
- Shows response variability across conditions

PANEL b - O2_dom_freq:
- Dominant frequency of O2 oscillations
- Captures rhythmic patterns in oxygen consumption

PANEL c - Amp_dom_freq:
- Dominant frequency of contractility amplitude
- Related to beating frequency/regularity

DATA SOURCE:
- Cleaned_Data/Heatmaps/Vandetanib (G11)/*.csv

GOAL: Show additional metrics beyond O2_mean and Amp_std.
"""
    },
}


def load_registry():
    """Load figure registry with sizes."""
    df = pd.read_csv(REGISTRY_PATH)
    df = df.drop_duplicates(subset=['Figure_ID', 'Letter'], keep='first')
    return df


def clear_slide_shapes(slide):
    """Remove all shapes from slide."""
    shapes_to_remove = list(slide.shapes)
    for shape in shapes_to_remove:
        sp = shape._element
        sp.getparent().remove(sp)


def add_title_to_slide(slide, title_text):
    """Add title to slide."""
    txBox = slide.shapes.add_textbox(Inches(MARGIN), Inches(MARGIN * 0.3),
                                      Inches(SLIDE_WIDTH - 2 * MARGIN), Inches(TITLE_HEIGHT))
    tf = txBox.text_frame
    p = tf.paragraphs[0]
    p.text = title_text
    p.font.size = Pt(12)
    p.font.bold = True
    p.alignment = PP_ALIGN.CENTER


def add_figure_with_label(slide, fig_path, left, top, width, height, label=None):
    """Add figure at exact size with optional label."""
    if not fig_path.exists():
        txBox = slide.shapes.add_textbox(Inches(left), Inches(top),
                                          Inches(width), Inches(height))
        tf = txBox.text_frame
        p = tf.paragraphs[0]
        p.text = f"[{fig_path.name}]"
        p.font.size = Pt(8)
        return

    slide.shapes.add_picture(str(fig_path), Inches(left), Inches(top),
                              width=Inches(width), height=Inches(height))

    if label:
        lbl_box = slide.shapes.add_textbox(Inches(left + 0.02), Inches(top + 0.02),
                                            Inches(0.2), Inches(0.2))
        tf = lbl_box.text_frame
        p = tf.paragraphs[0]
        p.text = label
        p.font.size = Pt(10)
        p.font.bold = True


def flow_layout_panels(panels, start_top):
    """Calculate positions for panels in a flow layout."""
    positions = []
    current_left = MARGIN
    current_top = start_top
    row_max_height = 0
    usable_width = SLIDE_WIDTH - 2 * MARGIN

    for width, height in panels:
        if current_left + width > MARGIN + usable_width and current_left > MARGIN:
            current_left = MARGIN
            current_top += row_max_height + GAP
            row_max_height = 0

        positions.append((current_left, current_top, width, height))
        current_left += width + GAP
        row_max_height = max(row_max_height, height)

    return positions


def populate_slide_from_registry(slide, fig_id, registry_df, config):
    """Populate a slide using registry data for sizes."""
    clear_slide_shapes(slide)

    # Add title
    add_title_to_slide(slide, config['title'])

    # Add notes
    if slide.has_notes_slide:
        notes_slide = slide.notes_slide
    else:
        notes_slide = slide.notes_slide

    notes_tf = notes_slide.notes_text_frame
    notes_tf.text = config['notes']

    # Get panels for this figure
    fig_panels = registry_df[registry_df['Figure_ID'] == fig_id].copy()

    if fig_panels.empty:
        return 0

    # Prepare panel info
    panels_info = []
    for _, row in fig_panels.iterrows():
        width = float(row['Width_In'])
        height = float(row['Height_In'])
        letter = row['Letter'] if pd.notna(row['Letter']) else ''
        png_path = PROJECT_ROOT / row['PNG_Path']
        panels_info.append((width, height, png_path, letter))

    # Calculate flow positions
    panel_sizes = [(w, h) for w, h, _, _ in panels_info]
    positions = flow_layout_panels(panel_sizes, MARGIN + TITLE_HEIGHT + 0.1)

    # Add each panel
    for (left, top, width, height), (_, _, png_path, letter) in zip(positions, panels_info):
        add_figure_with_label(slide, png_path, left, top, width, height, letter)

    return len(panels_info)


def main():
    registry_df = load_registry()
    registry_df['Figure_ID'] = registry_df['Figure_ID'].astype(str)

    if PPTX_PATH.exists():
        prs = Presentation(str(PPTX_PATH))
        print(f"Opened: {PPTX_PATH}")
    else:
        prs = Presentation()
        prs.slide_width = Inches(SLIDE_WIDTH)
        prs.slide_height = Inches(SLIDE_HEIGHT)
        print("Created new presentation")

    # Figure order
    fig_order = ['1', '2', '3', '4', '5', '6', '7', '8', 'S1']

    # Ensure enough slides
    while len(prs.slides) < len(fig_order):
        blank_layout = prs.slide_layouts[6]
        prs.slides.add_slide(blank_layout)

    # Populate each slide
    for i, fig_id in enumerate(fig_order):
        slide = prs.slides[i]
        config = FIGURE_CONFIGS[fig_id]
        n_panels = populate_slide_from_registry(slide, fig_id, registry_df, config)
        print(f"  Slide {i+1}: {config['title']} ({n_panels} panels) + notes")

    prs.save(str(PPTX_PATH))
    print(f"\nSaved: {PPTX_PATH}")
    print("Notes added to all slides with figure descriptions and data sources.")


if __name__ == '__main__':
    main()
