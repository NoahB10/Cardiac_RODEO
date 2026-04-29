"""Apply the new layout (Prism panels in their new slot mapping) to
`Output/PowerPoint_Figures_Remake/Cardiac_RODEO_Remake.pptx`.

Slide map:
    slide 7  -> Figure 6 (Arrhythmia)
    slide 8  -> Figure 7 (HeartDamage)
    slide 9  -> Figure 8 (ConcernBinary)

Slide 7 already has Panel_6{letter} groups in the right positions — we just
swap the picture INSIDE each group.  Slides 8 and 9 don't have groups, so
we rebuild their picture set from scratch using the same box positions.

Panel letters are placed/aligned per the same layout boxes.  No re-rendering
of Prism PNGs is performed here; if a Prism image's aspect ratio doesn't
match its box exactly, PPT's fit-to-box scaling will preserve aspect (the
image is centred inside the box, with whitespace at the unused dim).
"""

from __future__ import annotations

import json
import shutil
import sys
from pathlib import Path

from pptx import Presentation
from pptx.util import Emu, Pt, Inches
from pptx.enum.shapes import MSO_SHAPE_TYPE
from pptx.enum.text import PP_ALIGN
from pptx.dml.color import RGBColor
from lxml import etree

HERE = Path(__file__).resolve().parent
PROJECT_ROOT = HERE.parent
sys.path.insert(0, str(HERE))

from _layout import (ROW_LAYOUT, PANEL_ROW, MARGIN_B, CONTENT,
                     LEGEND_STASH_X, LEGEND_STASH_T_BY_LETTER,
                     LEGEND_FILE_BY_LETTER,
                     INPLACE_PANELS, INPLACE_FIG_NUM,
                     RESIZE_TO_NATIVE,
                     LOOSE_DESIGN_ANCHORS)  # noqa: E402
from _paths import panel_dir  # noqa: E402
from PIL import Image as _PILImage  # noqa: E402

# User-adjusted PanelText_* positions (captured by extract_text_overrides.py).
# Optional — if the file doesn't exist, sidecar auto-positions are used.
try:
    from _text_overrides import TEXT_OVERRIDES  # noqa: E402
except ImportError:
    TEXT_OVERRIDES = {}

# Pre-strip plot-area origins (panel-relative inches). Used to re-anchor the
# picture frame after STRIP_LABELS=True shrinks the PNG: we shift the frame
# so the plot area lands at the SAME slide coordinates as before. Without
# this, a stripped PNG would land at the original picture top-left and the
# plot area would visibly shift up/left.
try:
    from _plot_origins import PLOT_ORIGINS  # noqa: E402
except ImportError:
    PLOT_ORIGINS = {}

EMU_PER_INCH = 914400

PPTX_IN  = PROJECT_ROOT / "Output" / "PowerPoint_Figures_Remake" / "Cardiac_RODEO_Remake.pptx"
PPTX_OUT = PROJECT_ROOT / "Output" / "PowerPoint_Figures_Remake" / "Cardiac_RODEO_Remake.pptx"  # in-place

FIGS = (6, 7, 8)


def _find_figure_slides(prs) -> dict[int, list[int]]:
    """Find every slide whose first text matches 'Figure {N}:' or 'Figure {N} '.

    Returns {fig_num: [zero_based_slide_index, ...]} ordered by slide order.
    Re-runnable: doesn't depend on hardcoded slide indices.
    """
    out: dict[int, list[int]] = {}
    for i, sl in enumerate(prs.slides):
        title = ""
        for sp in sl.shapes:
            if sp.has_text_frame and sp.text_frame.text.strip():
                title = sp.text_frame.text.strip()
                break
        for fig in FIGS:
            if title.startswith(f"Figure {fig}:") or title.startswith(f"Figure {fig} "):
                out.setdefault(fig, []).append(i)
                break
    return out

# Letter placement: offsets relative to the box's top-left corner.
LETTER_OFFSET_X_IN = -0.05
LETTER_OFFSET_Y_IN = -0.18
LETTER_BOX_W_IN = 0.30
LETTER_BOX_H_IN = 0.20
LETTER_FONT_PT = 12
LETTER_FONT_BOLD = True

# Toggle: when False, the script does NOT add PanelLetter_* textboxes and
# strips any pre-existing ones from each managed slide (slides 2/3/6/7/8).
# Flip back to True to restore the panel-letter overlays. The _add_letter()
# function definition is kept intact regardless of this flag.
ADD_PANEL_LETTERS = True

# Toggle: when True, read each panel PNG's `.text.json` sidecar and add
# editable PowerPoint text boxes (named PanelText_{slide}{letter}_{role}_{idx})
# overlaying the rendered axis labels / titles / annotations. Stale boxes from
# prior runs are removed first. Set False to strip overlays without rebuilding.
ADD_PANEL_TEXT = True

# Per-slide opt-in: which slides currently get text overlays. We're rolling
# this out one slide at a time (start with Fig 2, then Fig 6/7/8, then Fig 3)
# so the user can verify alignment slide-by-slide before stripping labels
# from the PNGs. Add slide numbers (1-based) here as each is verified.
PANEL_TEXT_SLIDES = {2, 3, 6, 7, 8}

# Per-panel letter offset overrides (slide_1based, letter) -> (dx_in, dy_in).
# Applied RELATIVE to the picture's (L, T). Falls back to
# LETTER_OFFSET_X_IN / LETTER_OFFSET_Y_IN when not present.
# Slide 3 letters are row-aligned to Y = 0.79 (row 1 a-f), 1.87 (row 2 g/i),
# 3.93 (row 3 j/k). Heatmap panels (a/c/e) and the wide ROC panel (i) need
# negative dy because their pictures sit lower than the row-1 surface
# pictures; b/d/f keep their custom dx (artwork has whitespace at panel-left).
LETTER_OFFSETS: dict[tuple[int, str], tuple[float, float]] = {
    (3, "a"): (-0.05, -0.21),
    (3, "b"): (0.243, 0.124),
    (3, "c"): (-0.05, -0.28),
    (3, "d"): (0.309, 0.124),
    (3, "e"): (-0.05, -0.23),
    (3, "f"): (0.393, 0.124),
    (3, "i"): (-0.05, -0.32),
    (3, "j"): (-0.05, -0.19),
}

NS = {
    "a": "http://schemas.openxmlformats.org/drawingml/2006/main",
    "p": "http://schemas.openxmlformats.org/presentationml/2006/main",
    "r": "http://schemas.openxmlformats.org/officeDocument/2006/relationships",
}


def _swap_picture_source(sp, new_image_path: Path) -> None:
    """Replace the image bytes behind a python-pptx Picture, keeping its bbox."""
    a_ns = NS["a"]; r_ns = NS["r"]
    blip = sp._element.find(f".//{{{a_ns}}}blip")
    if blip is None:
        raise RuntimeError(f"no blip on shape {sp.name}")
    rid = blip.attrib.get(f"{{{r_ns}}}embed")
    rel = sp.part.rels[rid]
    rel.target_part._blob = Path(new_image_path).read_bytes()


def _delete_shape(slide, shape):
    sp = shape._element
    sp.getparent().remove(sp)


def _add_letter(slide, letter: str, position_in: tuple[float, float],
                slide_1based: int = 0):
    """Place a bold letter textbox near position_in.
    Uses LETTER_OFFSETS[(slide_1based, letter)] if defined, otherwise the
    global LETTER_OFFSET_X/Y_IN defaults."""
    L, T = position_in
    dx, dy = LETTER_OFFSETS.get((slide_1based, letter),
                                (LETTER_OFFSET_X_IN, LETTER_OFFSET_Y_IN))
    x = (L + dx) * EMU_PER_INCH
    y = (T + dy) * EMU_PER_INCH
    if x < 0:
        x = 0
    if y < 0:
        y = 0
    tb = slide.shapes.add_textbox(int(x), int(y),
                                  int(LETTER_BOX_W_IN * EMU_PER_INCH),
                                  int(LETTER_BOX_H_IN * EMU_PER_INCH))
    tb.name = f"PanelLetter_{letter}"
    tf = tb.text_frame
    tf.margin_left = tf.margin_right = tf.margin_top = tf.margin_bottom = 0
    p = tf.paragraphs[0]
    p.alignment = PP_ALIGN.LEFT
    run = p.add_run()
    run.text = letter
    run.font.name = "Arial"
    run.font.size = Pt(LETTER_FONT_PT)
    run.font.bold = LETTER_FONT_BOLD
    run.font.color.rgb = RGBColor(0, 0, 0)


# --------------------------------------------------------------------------- #
# Text overlay (sidecar JSON -> editable PowerPoint text boxes)
# --------------------------------------------------------------------------- #

_HALIGN_MAP = {
    "left": PP_ALIGN.LEFT, "center": PP_ALIGN.CENTER,
    "right": PP_ALIGN.RIGHT, "justify": PP_ALIGN.JUSTIFY,
}


def _read_text_sidecar(png_path: Path) -> dict | None:
    sidecar = Path(str(png_path) + ".text.json")
    if not sidecar.exists():
        return None
    try:
        return json.loads(sidecar.read_text())
    except Exception as e:
        print(f"  [WARN] failed to parse sidecar {sidecar.name}: {e}")
        return None


def _add_text_box(slide, element: dict, panel_LT_in: tuple[float, float],
                  shape_name: str, *, override: dict | None = None) -> None:
    """Create one PowerPoint text box.

    If ``override`` is given (from `_text_overrides.TEXT_OVERRIDES`), use it
    directly: it stores the panel-relative (x, y, w, h, rotation_deg) in
    PowerPoint's CW-positive convention — no unrotation math needed.

    Otherwise, fall back to the sidecar element's matplotlib bbox + rotation,
    and convert (rotated bbox → unrotated PPTX box, mpl CCW rot → pptx CW rot).
    """
    L, T = panel_LT_in

    if override is not None:
        x = L + override["x_in"]
        y = T + override["y_in"]
        w = override["w_in"]
        h = override["h_in"]
        rot_pptx = float(override.get("rotation_deg", 0.0) or 0.0)
    else:
        x = L + element["x_in"]
        y = T + element["y_in"]
        w = element["w_in"]
        h = element["h_in"]
        rot_mpl = element.get("rotation_deg", 0.0) or 0.0
        if abs(rot_mpl) > 1e-3:
            # Rotated text: swap dims and recenter so the rotated bounding
            # rect matches the captured (rotated) bbox. PowerPoint rotates
            # around the shape's center; the unrotated left may be negative
            # (extends off-slide pre-rotation) — that's fine.
            cx = x + w / 2
            cy = y + h / 2
            natural_w, natural_h = h, w   # for ±90°
            x = cx - natural_w / 2
            y = cy - natural_h / 2
            w, h = natural_w, natural_h
            # mpl CCW-positive -> pptx CW-positive.
            rot_pptx = -float(rot_mpl)
        else:
            rot_pptx = 0.0

    tb = slide.shapes.add_textbox(int(x * EMU_PER_INCH),
                                  int(y * EMU_PER_INCH),
                                  int(w * EMU_PER_INCH),
                                  int(h * EMU_PER_INCH))
    tb.name = shape_name

    if abs(rot_pptx) > 1e-3:
        tb.rotation = rot_pptx

    tf = tb.text_frame
    tf.margin_left = tf.margin_right = tf.margin_top = tf.margin_bottom = 0
    tf.word_wrap = False

    p = tf.paragraphs[0]
    p.alignment = _HALIGN_MAP.get(element.get("halign", "left"), PP_ALIGN.LEFT)

    run = p.add_run()
    run.text = element["text"]
    run.font.name = "Arial"
    run.font.size = Pt(float(element.get("font_pt", 13)))
    run.font.bold = bool(element.get("bold", False))
    rgb = element.get("color_rgb") or [0, 0, 0]
    run.font.color.rgb = RGBColor(int(rgb[0]), int(rgb[1]), int(rgb[2]))


def _walk_shapes_with_parent(shapes, parent=None):
    """Yield (shape, container_shapes) descending into groups.

    container_shapes is the SpTree-like container that owns the shape — used
    by _delete_shape's getparent() path. Needed when the user has manually
    grouped PanelText boxes inside Panel_X groups.
    """
    for sp in shapes:
        if sp.shape_type == MSO_SHAPE_TYPE.GROUP:
            yield from _walk_shapes_with_parent(sp.shapes, sp)
        else:
            yield sp, parent


def _remove_panel_text_for(slide, slide_1based: int, letter: str,
                            sidecar_keys: set | None = None) -> int:
    """Remove PanelText_{slide}{letter}_* shapes from the slide.

    If ``sidecar_keys`` is provided, only shapes whose (role, idx) appears
    in the set are removed — "orphan" shapes (overrides for elements that
    don't appear in the current sidecar, e.g. an axis label baked-out by
    STRIP_LABELS) stay in place. Pass None to remove all matching shapes.

    Descends into groups (the user may have manually grouped overlay text
    boxes inside their parent panel group in PowerPoint).
    """
    prefix = f"PanelText_{slide_1based}{letter}_"
    removed = 0
    for shape, _parent in list(_walk_shapes_with_parent(slide.shapes)):
        if not (hasattr(shape, "name") and shape.name.startswith(prefix)):
            continue
        if sidecar_keys is not None:
            rest = shape.name[len(prefix):]
            head, _, tail = rest.rpartition("_")
            try:
                idx = int(tail)
            except ValueError:
                continue
            if (head, idx) not in sidecar_keys:
                continue
        _delete_shape(slide, shape)
        removed += 1
    return removed


def _remove_all_panel_text(slide) -> int:
    """Remove every PanelText_* shape from a slide (used when ADD_PANEL_TEXT=False)."""
    removed = 0
    for shape, _parent in list(_walk_shapes_with_parent(slide.shapes)):
        if hasattr(shape, "name") and shape.name.startswith("PanelText_"):
            _delete_shape(slide, shape)
            removed += 1
    return removed


def _panel_override_keys(slide_1based: int, letter: str) -> set:
    """Return the set of (role, idx) keys in TEXT_OVERRIDES for this panel."""
    return {(role, idx) for (s, l, role, idx) in TEXT_OVERRIDES
            if s == slide_1based and l == letter}


def _ensure_orphan_overlays(slide, panel_LT_in, slide_1based, letter,
                            sidecar_keys: set) -> int:
    """Create PanelText shapes for TEXT_OVERRIDES entries that have no
    matching sidecar element (e.g. axis labels stripped by STRIP_LABELS).

    Skips entries whose shape already exists on the slide. Uses the override
    dict's saved 'text' / 'font_pt' / 'bold' / 'color_rgb' fields.
    """
    panel_keys = _panel_override_keys(slide_1based, letter)
    orphan_keys = panel_keys - sidecar_keys
    if not orphan_keys:
        return 0

    existing = {sp.name for sp, _ in _walk_shapes_with_parent(slide.shapes)
                if hasattr(sp, "name")}

    created = 0
    for role, idx in orphan_keys:
        shape_name = f"PanelText_{slide_1based}{letter}_{role}_{idx}"
        if shape_name in existing:
            continue
        override = TEXT_OVERRIDES[(slide_1based, letter, role, idx)]
        # Build a minimal "element" dict from the override's saved props.
        element = {
            "role": role,
            "text": override.get("text", ""),
            "font_pt": override.get("font_pt", 13.0),
            "bold": override.get("bold", False),
            "color_rgb": override.get("color_rgb", [0, 0, 0]),
            "halign": "left",
        }
        _add_text_box(slide, element, panel_LT_in, shape_name, override=override)
        created += 1
    return created


def _add_panel_text_overlays(slide, png_path: Path,
                              panel_LT_in: tuple[float, float],
                              slide_1based: int, letter: str) -> int:
    """Add PanelText overlays for this panel.

    The panel is in CURATED mode if TEXT_OVERRIDES has any entry for it. In
    that mode, only sidecar elements whose (role, idx) match an override key
    get an overlay box — the rest are assumed BAKED in the PNG (the user
    has explicitly opted out of overlays for them, e.g. CM cell numbers,
    bar values). Orphan overrides (overrides without sidecar match) are
    created from the override's saved text+font props.

    In BOOTSTRAP mode (no overrides exist for the panel), every sidecar
    element gets an overlay at its auto-position — this is the first-pass
    behavior used to seed text boxes that the user then adjusts.

    Returns the count of overlay boxes added/replaced.
    """
    sidecar = _read_text_sidecar(png_path)
    sidecar_elements = (sidecar or {}).get("elements") or []
    sidecar_keys = {(el.get("role", "text"), idx)
                    for idx, el in enumerate(sidecar_elements)}

    panel_keys = _panel_override_keys(slide_1based, letter)
    curated = bool(panel_keys)

    # Remove only sidecar-matching shapes that we'll re-add below.
    removed_keys = set()
    if curated:
        # Curated mode — re-add only intersection of sidecar + overrides.
        removed_keys = sidecar_keys & panel_keys
    else:
        # Bootstrap mode — re-add every sidecar element.
        removed_keys = sidecar_keys
    if removed_keys:
        _remove_panel_text_for(slide, slide_1based, letter,
                               sidecar_keys=removed_keys)

    # Recreate overlays for sidecar elements (curated: filtered by overrides).
    overrides_used = 0
    created = 0
    for idx, el in enumerate(sidecar_elements):
        role = el.get("role", "text")
        if curated and (role, idx) not in panel_keys:
            continue   # user opted out — leave label baked in PNG
        shape_name = f"PanelText_{slide_1based}{letter}_{role}_{idx}"
        override = TEXT_OVERRIDES.get((slide_1based, letter, role, idx))
        if override is not None:
            overrides_used += 1
        _add_text_box(slide, el, panel_LT_in, shape_name, override=override)
        created += 1

    # Create orphan overlays (in TEXT_OVERRIDES but not in sidecar).
    orphans = _ensure_orphan_overlays(slide, panel_LT_in,
                                      slide_1based, letter, sidecar_keys)

    if overrides_used or orphans:
        msg = f"    text overrides applied: {overrides_used}/{created}"
        if orphans:
            msg += f"  +{orphans} orphan(s)"
        print(msg)
    return created + orphans


def _native_size_in(image_path: Path) -> tuple[float, float]:
    """Read the PNG's native size in INCHES from its embedded DPI metadata."""
    im = _PILImage.open(image_path)
    dpi = im.info.get("dpi", (600, 600))[0]
    return im.size[0] / dpi, im.size[1] / dpi


def _add_picture(slide, image_path: Path, position_in: tuple[float, float]):
    """Place picture at native size — no width/height override = no stretching."""
    L, T = position_in
    w_in, h_in = _native_size_in(image_path)
    return slide.shapes.add_picture(
        str(image_path),
        Inches(L), Inches(T),
        width=Inches(w_in), height=Inches(h_in),
    )


def _is_any_picture(shape):
    """Any picture shape — used to wipe everything image-shaped on a figure
    slide before rebuilding.  Catches small legend PNGs that fall below the
    old 0.5" 'panel picture' threshold."""
    return shape.shape_type == MSO_SHAPE_TYPE.PICTURE


def _is_panel_letter(shape):
    if not (shape.has_text_frame and shape.text_frame.text.strip()):
        return False
    txt = shape.text_frame.text.strip()
    return (len(txt) == 1 and txt.lower() in "abcdefghijk"
            and shape.name.startswith(("PanelLetter_", "Label_")))


def _remove_existing_panel_letters(slide) -> int:
    """Strip any pre-existing PanelLetter_* / Label_* textboxes from a slide.
    Used when ADD_PANEL_LETTERS is False so re-runs don't leave stale labels.
    Returns the count removed."""
    removed = 0
    for shape in list(slide.shapes):
        if _is_panel_letter(shape):
            _delete_shape(slide, shape)
            removed += 1
    return removed


def _compute_picture_top(letter: str, image_h_in: float) -> float:
    """Top position so this panel's plot bottom sits at its row's plot_bottom.

    plot_bottom = T + image_h - margin_b   =>   T = plot_bottom + margin_b - image_h
    """
    row = ROW_LAYOUT[PANEL_ROW[letter]]
    return row["plot_bottom"] + MARGIN_B[letter] - image_h_in


def update_slide_loose(slide, fig_num: int, slide_1based: int = 0):
    """Wipe ALL pictures + panel letters + groups, then rebuild from
    ROW_LAYOUT + CONTENT. Title text boxes are preserved.

    Removing ALL pictures (not just 'large' ones) prevents legend PNGs from
    accumulating across re-runs.
    """
    for shape in list(slide.shapes):
        if shape.shape_type == MSO_SHAPE_TYPE.GROUP:
            _delete_shape(slide, shape)
    for shape in list(slide.shapes):
        if _is_any_picture(shape):
            _delete_shape(slide, shape)
        elif _is_panel_letter(shape):
            _delete_shape(slide, shape)

    panel_positions: dict[str, tuple[float, float, float, float]] = {}
    for name in sorted(CONTENT[fig_num].keys()):
        letter = name[-1]
        png_name = CONTENT[fig_num][name]
        png = panel_dir(fig_num) / png_name
        if not png.exists():
            print(f"  [WARN] missing source {png.name} for {name}")
            continue
        w_in, h_in = _native_size_in(png)
        row = ROW_LAYOUT[PANEL_ROW[letter]]
        L_baseline = row["lefts"][letter]
        T_baseline = _compute_picture_top(letter, h_in)

        # Text-box anchor (= design-time picture top-left). Stays invariant
        # across STRIP_LABELS modes so user-adjusted overrides round-trip.
        # In baked mode (panel not yet snapshotted), anchor = the freshly
        # computed (L, T). In strip mode, anchor = LOOSE_DESIGN_ANCHORS
        # snapshot taken before the flip.
        anchor_L, anchor_T = LOOSE_DESIGN_ANCHORS.get(
            (slide_1based, letter), (L_baseline, T_baseline))

        # Picture position. If the panel is opted into the text-overlay
        # system AND has a PLOT_ORIGINS entry AND the sidecar is available,
        # anchor the picture so the plot area lands at PLOT_ORIGINS.
        plot_origin = PLOT_ORIGINS.get((slide_1based, letter))
        sidecar = _read_text_sidecar(png)
        sidecar_plot = (sidecar or {}).get("plot_area")
        anchor_to_plot = (slide_1based in PANEL_TEXT_SLIDES
                          and plot_origin is not None
                          and sidecar_plot is not None)
        if anchor_to_plot:
            L = plot_origin[0] - sidecar_plot["x_in"]
            T = plot_origin[1] - sidecar_plot["y_in"]
        else:
            L, T = L_baseline, T_baseline

        _add_picture(slide, png, (L, T))
        if ADD_PANEL_LETTERS:
            _add_letter(slide, letter, (L, row["letter_top"] - LETTER_OFFSET_Y_IN),
                        slide_1based=slide_1based)
        panel_positions[letter] = (L, T, w_in, h_in)
        plot_bottom = T + h_in - MARGIN_B[letter]
        if anchor_to_plot:
            shift_x = L - L_baseline
            shift_y = T - T_baseline
            print(f"  add {name:12s} <- {png_name}  pos=({L:.2f},{T:.2f}) "
                  f"native={w_in:.2f}x{h_in:.2f}\"  "
                  f"shift=({shift_x:+.2f},{shift_y:+.2f})\"")
        else:
            print(f"  add {name:12s} <- {png_name}  pos=({L:.2f},{T:.2f}) "
                  f"native={w_in:.2f}x{h_in:.2f}\"  plot_bottom={plot_bottom:.2f}")

        if ADD_PANEL_TEXT and slide_1based in PANEL_TEXT_SLIDES:
            count = _add_panel_text_overlays(
                slide, png, (anchor_L, anchor_T), slide_1based, letter)
            if count:
                print(f"    text overlays: {count} from sidecar "
                      f"(anchor=({anchor_L:.2f},{anchor_T:.2f}))")

    # Legend overlays — stashed in the off-slide grey area (L > slide width)
    # so the user can drag each one onto the panel manually.  Only one of
    # each legend per slide; on re-runs the wipe step above clears any
    # prior copies.
    for legend_letter, suffix in LEGEND_FILE_BY_LETTER.items():
        if legend_letter not in panel_positions:
            continue
        legend_png = panel_dir(fig_num) / f"Fig_{fig_num}{suffix}"
        if not legend_png.exists():
            print(f"  [SKIP legend] missing {legend_png.name}")
            continue
        legend_L = LEGEND_STASH_X
        legend_T = LEGEND_STASH_T_BY_LETTER[legend_letter]
        _add_picture(slide, legend_png, (legend_L, legend_T))
        leg_w, leg_h = _native_size_in(legend_png)
        print(f"  stash legend <- {legend_png.name}  "
              f"pos=({legend_L:.2f},{legend_T:.2f}) (off-slide) "
              f"native={leg_w:.2f}x{leg_h:.2f}\"")


def _delete_slide(prs, slide_idx_zero_based: int):
    """Remove a slide from the presentation by 0-based index."""
    rId = prs.slides._sldIdLst[slide_idx_zero_based].rId
    prs.part.drop_rel(rId)
    del prs.slides._sldIdLst[slide_idx_zero_based]


def _walk_pictures(shapes):
    """Yield every picture shape, descending into groups."""
    for sp in shapes:
        if sp.shape_type == MSO_SHAPE_TYPE.GROUP:
            yield from _walk_pictures(sp.shapes)
        elif sp.shape_type == MSO_SHAPE_TYPE.PICTURE:
            yield sp


def update_inplace_panels(prs, *, position_tol_in: float = 0.05):
    """Swap source bytes for in-place panel pictures on slides 2 and 3.

    Covers heatmaps AND line/sigmoid background panels — anything where the
    user has already placed the picture frame manually. The picture frames
    are identified by their (left_in, top_in) position; we don't move or
    resize them. Works for free-standing pictures AND for the BACKGROUND
    picture inside a group (whose left/top sits at the group origin).

    If a frame at the expected position isn't found, we log a warning and
    skip — the user may have repositioned it, in which case the swap should
    be redone manually.
    """
    for (slide_1based, letter), (exp_L, exp_T, png_name) in INPLACE_PANELS.items():
        if slide_1based - 1 >= len(prs.slides):
            print(f"  [SKIP] slide {slide_1based} out of range "
                  f"(deck has {len(prs.slides)})")
            continue
        slide = prs.slides[slide_1based - 1]
        fig_num = INPLACE_FIG_NUM[slide_1based]
        png_path = panel_dir(fig_num) / png_name
        if not png_path.exists():
            print(f"  [SKIP] {png_name} not found at {png_path}")
            continue

        # Look for the picture frame. For non-anchored panels, the frame must
        # be at INPLACE spec ± 0.05". For anchored panels (slide in
        # PANEL_TEXT_SLIDES + PLOT_ORIGINS available), the frame may have
        # been shifted by a previous apply with different strip-mode margins,
        # so accept anything in the rectangle from the spec position to
        # spec + max_possible_shift. Max shift = (PLOT_ORIGIN - INPLACE) in
        # both axes (= the original PNG margin), since stripped margin is
        # always smaller.
        plot_origin = PLOT_ORIGINS.get((slide_1based, letter))
        anchored_panel = (slide_1based in PANEL_TEXT_SLIDES
                          and plot_origin is not None)
        match = None
        if anchored_panel:
            max_shift_x = max(0.0, plot_origin[0] - exp_L)
            max_shift_y = max(0.0, plot_origin[1] - exp_T)
            tol_x = max(position_tol_in, max_shift_x + 0.05)
            tol_y = max(position_tol_in, max_shift_y + 0.05)
            # Frame's left/top must be ≥ exp - 0.05 and ≤ exp + tol.
            for sp in _walk_pictures(slide.shapes):
                L = sp.left / EMU_PER_INCH
                T = sp.top / EMU_PER_INCH
                if (exp_L - 0.05 <= L <= exp_L + tol_x
                        and exp_T - 0.05 <= T <= exp_T + tol_y):
                    match = sp
                    break
        else:
            for sp in _walk_pictures(slide.shapes):
                L = sp.left / EMU_PER_INCH
                T = sp.top / EMU_PER_INCH
                if (abs(L - exp_L) <= position_tol_in
                        and abs(T - exp_T) <= position_tol_in):
                    match = sp
                    break

        if match is None:
            print(f"  [WARN] slide {slide_1based} panel {letter}: no picture "
                  f"frame near ({exp_L:.2f}, {exp_T:.2f})")
            continue

        _swap_picture_source(match, png_path)

        # If this panel is in PANEL_TEXT_SLIDES + PLOT_ORIGINS, anchor the
        # picture so the plot area lands at INPLACE_PANELS coords + saved
        # plot origin. Required when STRIP_LABELS shrinks the PNG: without
        # the shift, the plot area would visibly move up/left because the
        # smaller PNG's left/top margin no longer matches the original.
        plot_origin = PLOT_ORIGINS.get((slide_1based, letter))
        sidecar = _read_text_sidecar(png_path)
        sidecar_plot = (sidecar or {}).get("plot_area")
        anchor_to_plot = (slide_1based in PANEL_TEXT_SLIDES
                          and plot_origin is not None
                          and sidecar_plot is not None)

        if anchor_to_plot:
            new_w_in, new_h_in = _native_size_in(png_path)
            # PLOT_ORIGINS now stores absolute slide-coord plot top-left.
            target_plot_x = plot_origin[0]
            target_plot_y = plot_origin[1]
            new_pic_L = target_plot_x - sidecar_plot["x_in"]
            new_pic_T = target_plot_y - sidecar_plot["y_in"]
            match.left = Inches(new_pic_L)
            match.top = Inches(new_pic_T)
            match.width = Inches(new_w_in)
            match.height = Inches(new_h_in)
            shift_x = new_pic_L - exp_L
            shift_y = new_pic_T - exp_T
            print(f"  swap+anchor slide {slide_1based} panel {letter:s} <- "
                  f"{png_name}  frame=({new_pic_L:.2f},{new_pic_T:.2f}) "
                  f"{new_w_in:.2f}x{new_h_in:.2f}\"  "
                  f"shift=({shift_x:+.2f},{shift_y:+.2f})\"")
        elif (slide_1based, letter) in RESIZE_TO_NATIVE:
            new_w_in, new_h_in = _native_size_in(png_path)
            match.width = Inches(new_w_in)
            match.height = Inches(new_h_in)
            print(f"  swap+resize slide {slide_1based} panel {letter:s} <- "
                  f"{png_name}  frame=({exp_L:.2f},{exp_T:.2f}) "
                  f"{new_w_in:.2f}x{new_h_in:.2f}\" (native)")
        else:
            w_in = match.width / EMU_PER_INCH
            h_in = match.height / EMU_PER_INCH
            print(f"  swap slide {slide_1based} panel {letter:s} <- {png_name}  "
                  f"frame=({exp_L:.2f},{exp_T:.2f}) {w_in:.2f}x{h_in:.2f}\"")

        # Picture-letter anchor uses the picture's CURRENT position (which may
        # have just shifted under anchor-to-plot). Text boxes use the
        # INPLACE_PANELS SPEC L,T as a stable anchor — independent of any
        # picture shift — so user-placed text boxes don't move when the
        # picture shifts. Spec ≈ pre-shift frame position within ±0.005".
        L_fr = match.left / EMU_PER_INCH
        T_fr = match.top / EMU_PER_INCH
        L_anchor = exp_L
        T_anchor = exp_T

        if ADD_PANEL_LETTERS:
            # Remove any stale PanelLetter for this slot before adding a fresh one.
            for sp in list(slide.shapes):
                if hasattr(sp, 'name') and sp.name == f"PanelLetter_{letter}":
                    _delete_shape(slide, sp)
            _add_letter(slide, letter, (L_fr, T_fr), slide_1based=slide_1based)
            dx, dy = LETTER_OFFSETS.get((slide_1based, letter),
                                        (LETTER_OFFSET_X_IN, LETTER_OFFSET_Y_IN))
            print(f"    letter '{letter}' at ({L_fr + dx:.2f}, {T_fr + dy:.2f}\")")

        if ADD_PANEL_TEXT and slide_1based in PANEL_TEXT_SLIDES:
            count = _add_panel_text_overlays(
                slide, png_path, (L_anchor, T_anchor), slide_1based, letter)
            if count:
                print(f"    text overlays: {count} from sidecar")
        else:
            # Slide isn't opted in yet — sweep any prior PanelText_* for this
            # slot so re-runs don't leave stale boxes behind.
            _remove_panel_text_for(slide, slide_1based, letter)


def main():
    PPTX_OUT.parent.mkdir(parents=True, exist_ok=True)
    if PPTX_IN != PPTX_OUT:
        shutil.copy(PPTX_IN, PPTX_OUT)
    prs = Presentation(str(PPTX_OUT))

    found = _find_figure_slides(prs)
    canonical = {}      # fig_num -> 0-based slide index of the slide we'll keep
    duplicates = []     # 0-based indices to delete (older copies of each fig)
    for fig, idxs in found.items():
        # Prefer the LAST occurrence — that's where any recent edits live
        # (e.g. the user's manual alignment work).
        canonical[fig] = idxs[-1]
        duplicates.extend(idxs[:-1])

    # Phase 1: rebuild canonical slides with plot-base alignment.
    for fig, idx in canonical.items():
        slide = prs.slides[idx]
        print(f"\n=== Slide {idx + 1} (Figure {fig}) ===")
        update_slide_loose(slide, fig, slide_1based=idx + 1)

    # Phase 2: drop duplicate slides AFTER processing so canonical indices
    # don't shift mid-loop. Delete from highest-index first.
    if duplicates:
        print(f"\n=== Removing duplicates: slides "
              f"{[i + 1 for i in duplicates]} ===")
        for idx in sorted(duplicates, reverse=True):
            _delete_slide(prs, idx)
            print(f"  removed slide {idx + 1}")

    # Phase 3: swap in-place panel picture bytes on slides 2 and 3
    # (Figures 2 and 3). Frames stay where the user placed them; the picture
    # bytes get refreshed from the latest Prism PNG.
    print("\n=== In-place panels (slides 2, 3) ===")
    update_inplace_panels(prs)

    # Phase 4: when panel-letter overlays are disabled, sweep every managed
    # slide and remove any pre-existing PanelLetter_* / Label_* textboxes so
    # re-runs don't leave stale labels behind.
    if not ADD_PANEL_LETTERS:
        print("\n=== Stripping panel-letter overlays (ADD_PANEL_LETTERS=False) ===")
        managed_slide_indices = set()
        # In-place slides 2 and 3 (1-based -> 0-based)
        for slide_1based in INPLACE_FIG_NUM:
            if slide_1based - 1 < len(prs.slides):
                managed_slide_indices.add(slide_1based - 1)
        # Loose-rebuild figure slides 6/7/8
        managed_slide_indices.update(canonical.values())
        for idx in sorted(managed_slide_indices):
            removed = _remove_existing_panel_letters(prs.slides[idx])
            if removed:
                print(f"  slide {idx + 1}: removed {removed} panel-letter textbox(es)")

    # Phase 5: same idea for the text-overlay system. When ADD_PANEL_TEXT is
    # off, sweep PanelText_* off every managed slide.
    if not ADD_PANEL_TEXT:
        print("\n=== Stripping panel-text overlays (ADD_PANEL_TEXT=False) ===")
        managed_slide_indices = set()
        for slide_1based in INPLACE_FIG_NUM:
            if slide_1based - 1 < len(prs.slides):
                managed_slide_indices.add(slide_1based - 1)
        managed_slide_indices.update(canonical.values())
        for idx in sorted(managed_slide_indices):
            removed = _remove_all_panel_text(prs.slides[idx])
            if removed:
                print(f"  slide {idx + 1}: removed {removed} panel-text textbox(es)")

    prs.save(str(PPTX_OUT))
    print(f"\n[done] -> {PPTX_OUT}")

    # Phase 6: re-group panels on slides 6/7/8 into Panel_{slide}{letter}
    # groups (and prune any extras + keep slide-3 surface labels in sync).
    # update_slide_loose blanket-deletes groups on each apply, so the
    # grouping must be re-applied after every save.
    try:
        import importlib
        import finalize_panels  # noqa: E402
        importlib.reload(finalize_panels)
        print("\n=== Finalize: re-group panels + sync surfaces ===")
        finalize_panels.main()
    except Exception as e:
        print(f"\n[WARN] finalize_panels failed: {e}")


if __name__ == "__main__":
    main()
